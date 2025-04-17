import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

#! styling the background, just to make it look more fasinating
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_layout_styled_background(prs, slide, layout_type):
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_background_fill(color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*color)
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def add_center_card(color, opacity=200):
        card_width = int(slide_width)
        card_height = int(slide_height * 0.75)
        card_left = int((slide_width - card_width) / 2)
        card_top = int((slide_height - card_height) / 1)

        card = slide.shapes.add_shape(
            # MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height
            MSO_SHAPE.RECTANGLE, card_left, card_top, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*color)
        card.fill.fore_color.alpha = opacity
        card.line.fill.background()
        slide.shapes._spTree.remove(card._element)
        slide.shapes._spTree.insert(3, card._element)

    def add_bottom_stripe(color, height_ratio=0.1):
        stripe_height = int(slide_height * height_ratio)
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, slide_height - stripe_height, slide_width, stripe_height
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(*color)
        stripe.line.fill.background()
        slide.shapes._spTree.remove(stripe._element)
        slide.shapes._spTree.insert(3, stripe._element)

    def add_side_cards(left_color, right_color, opacity=180):
        card_width = int(slide_width * 0.42)
        card_height = int(slide_height * 0.7)
        padding = int(slide_width * 0.04)
        top = int(slide_height * 0.15)

        for i, color in enumerate([left_color, right_color]):
            left = padding + i * (card_width + padding)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*color)
            card.fill.fore_color.alpha = opacity
            card.line.fill.background()
            slide.shapes._spTree.remove(card._element)
            slide.shapes._spTree.insert(3, card._element)

    # Style each layout uniquely
    if layout_type == "title_slide":
        add_background_fill((39, 68, 114))  # light_color
        add_center_card((234, 239, 242), 200)

    elif layout_type == "title_only":
        add_background_fill((234, 239, 242))

    elif layout_type == "title_and_content":
        add_background_fill((234, 239, 242))  # Slightly lighter
        add_bottom_stripe((39, 68, 114)) # dark color here

    elif layout_type == "section_header":
        add_background_fill((234, 239, 242))
        top_half = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height // 2)
        top_half.fill.solid()
        top_half.fill.fore_color.rgb = RGBColor(39, 68, 114)
        top_half.line.fill.background()
        slide.shapes._spTree.remove(top_half._element)
        slide.shapes._spTree.insert(3, top_half._element)

    elif layout_type == "two_content":
        add_background_fill((39, 68, 114))
        add_side_cards((234, 239, 242), (234, 239, 242), opacity=180)

    elif layout_type == "comparison":
        add_background_fill((39, 68, 114))
        add_side_cards((234, 239, 242), (234, 239, 242), opacity=180)

    elif layout_type in ["content_with_caption", "image_with_caption"]:
        add_background_fill((234, 239, 242))

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            int(slide_height * 0.15),
            int(slide_width * 0.02),
            int(slide_height * 0.7)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(39, 68, 114)
        bar.line.fill.background()
        slide.shapes._spTree.remove(bar._element)
        slide.shapes._spTree.insert(3, bar._element)

        # 🔧 Move and resize text box
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text.strip() != "":
                shape.left = int(slide_width * 0.06)
                shape.top = int(slide_height * 0.2)
                shape.width = int(slide_width * 0.5)
                shape.height = int(slide_height * 0.6)

        # 🔧 Move and resize image box
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape.left = int(slide_width * 0.6)
                shape.top = int(slide_height * 0.25)
                shape.width = int(slide_width * 0.35)
                shape.height = int(slide_height * 0.5)


# Configuration constants
MAX_BULLET_POINTS_PER_SLIDE = 8  # Maximum number of bullet points per slide
MAX_CONTENT_HEIGHT = 5  # Maximum content height in inches

#helper function
def set_bullet_points(text_frame, content, max_line_length=80, font_size=20):
    text_frame.clear()
    
    # First split by explicit newlines in the content
    paragraphs = content.split('\n')
    
    for i, paragraph in enumerate(paragraphs):
        # Add a new paragraph for each explicit newline in the input
        para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        
        # Format text to respect line length but keep all text in the same paragraph
        formatted_text = format_text(paragraph, max_line_length)
        para.text = formatted_text
        para.level = 0
        para.font.size = Pt(font_size)
    
    return len(paragraphs)  # Return the number of bullet points added

def check_content_overflow(paragraphs, max_bullets=MAX_BULLET_POINTS_PER_SLIDE):
    """Check if content will overflow the slide and split if needed."""
    if len(paragraphs) <= max_bullets:
        return paragraphs, None
    
    # Split content
    first_part = paragraphs[:max_bullets]
    overflow_part = paragraphs[max_bullets:]
    
    return first_part, overflow_part

def add_title_only_slide(prs, slide_data):
    """Add a slide with only a title."""
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_only")
    
    title = slide.shapes.title
    title.text = slide_data["title"]
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(44)
    title_format.bold = True
    
    return slide

def add_title_slide(prs, slide_data):
    """Add a title slide with title and subtitle."""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_slide")
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Set title text
    title.text = slide_data["title"]
    
    # Apply formatting to ALL paragraphs in the title text frame
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
    
    # Check which key exists in the slide_data
    sub_heading_key = "sub-heading" if "sub-heading" in slide_data else "sub_heading"
    subtitle.text = slide_data.get(sub_heading_key, "")
    
    # Apply formatting to ALL paragraphs in the subtitle text frame
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.italic = True
    
    return slide

# def add_title_and_content_slide(prs, slide_data, is_continuation=False):
#     """Add a slide with title and content with overflow handling."""
#     slide_layout = prs.slide_layouts[1]  # Title and Content
#     slide = prs.slides.add_slide(slide_layout)
#     add_layout_styled_background (prs, slide, layout_type="title_and_content")
    
#     title = slide.shapes.title
#     content = slide.placeholders[1]
    
#     # Adjust title for continuation slides
#     if is_continuation:
#         title.text = f"{slide_data['title']} (cont.)"
#     else:
#         title.text = slide_data["title"]
    
#     # Format title
#     title_format = title.text_frame.paragraphs[0].font
#     title_format.size = Pt(36)
#     title_format.bold = True
    
#     # Split content by newlines for overflow detection
#     content_paragraphs = slide_data["content"].split('\n')
    
#     # Check for overflow
#     first_part, overflow = check_content_overflow(content_paragraphs)
    
#     # Add first part content
#     first_content = '\n'.join(first_part)
#     set_bullet_points(content.text_frame, first_content, 80, 20)
    
#     # Handle overflow if present
#     if overflow:
#         # Create continuation slide with remaining content
#         continuation_data = slide_data.copy()
#         continuation_data["content"] = '\n'.join(overflow)
#         add_title_and_content_slide(prs, continuation_data, True)
    
#     return slide
def add_title_and_content_slide(prs, slide_data, is_continuation=False):
    """Add a slide with title and content with precise vertical overflow handling."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_and_content")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
    
    # Format all paragraphs in the title text frame
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True

    # Calculate the available height in the content placeholder
    available_height = content.height
    
    # Split content by newlines
    content_paragraphs = slide_data["content"].split('\n')
    
    # Create a testing function to determine exactly how many paragraphs will fit
    def test_content_fit(paragraphs, font_size=20):
        # Create a temporary slide for testing
        temp_slide = prs.slides.add_slide(prs.slide_layouts[1])
        temp_content = temp_slide.placeholders[1]
        
        # Set the content
        test_content = '\n'.join(paragraphs)
        set_bullet_points(temp_content.text_frame, test_content, 80, font_size)
        
        # Check if content fits within available height
        text_height = 0
        for p in temp_content.text_frame.paragraphs:
            # Get line spacing safely, default to 1.0 if None
            line_spacing = 1.0  # Default spacing
            if hasattr(p, 'line_spacing') and p.line_spacing is not None:
                line_spacing = p.line_spacing
                
            # Convert font size to Pt if it's not already
            font_size_pt = font_size
            if not isinstance(font_size_pt, Pt):
                font_size_pt = Pt(font_size)
                
            # Calculate paragraph height safely
            para_height = font_size_pt.pt * line_spacing * 1.2  # Use .pt attribute to get the numeric value
            text_height += para_height
        
        # Remove the temporary slide
        prs.slides._sldIdLst.remove(temp_slide._element.sldId)
        
        return text_height <= available_height
    
    # Use a simpler approach to avoid errors with temporary slides
    # For most presentations, using a fixed number based on paragraph length works well
    
    # Calculate average paragraph length
    avg_chars = sum(len(p) for p in content_paragraphs) / max(1, len(content_paragraphs))
    
    # Determine paragraphs per slide based on length
    if avg_chars > 200:  # Very long paragraphs
        max_paragraphs = 3
    elif avg_chars > 100:  # Medium paragraphs
        max_paragraphs = 4
    elif avg_chars > 50:  # Short paragraphs
        max_paragraphs = 6
    else:  # Very short paragraphs
        max_paragraphs = 8
        
    # Be more conservative for continuation slides
    if is_continuation:
        max_paragraphs = max(1, max_paragraphs - 1)
        
    # Calculate how many paragraphs to use on this slide
    paragraphs_to_use = min(max_paragraphs, len(content_paragraphs))
    
    # Split content
    first_part = content_paragraphs[:paragraphs_to_use]
    overflow = content_paragraphs[paragraphs_to_use:] if paragraphs_to_use < len(content_paragraphs) else None
    
    # Add first part content
    first_content = '\n'.join(first_part)
    set_bullet_points(content.text_frame, first_content, 80, 20)
    
    # Handle overflow if present
    if overflow and len(overflow) > 0:
        # Create continuation slide with remaining content
        continuation_data = slide_data.copy()
        continuation_data["content"] = '\n'.join(overflow)
        add_title_and_content_slide(prs, continuation_data, True)
    
    return slide

def add_two_content_slide(prs, slide_data, is_continuation=False):
    """Add a slide with title and two content areas with overflow handling."""
    slide_layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="two_content")

    
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(36)
    title_format.bold = True
    
    # Split content by newlines
    all_paragraphs = slide_data["content"].split('\n')
    total_paragraphs = len(all_paragraphs)
    
    # For this layout, we'll display half the content on each side
    max_per_side = MAX_BULLET_POINTS_PER_SLIDE // 2
    
    # Check if total content will fit
    if total_paragraphs <= MAX_BULLET_POINTS_PER_SLIDE:
        # Content will fit, split between columns
        mid_point = total_paragraphs // 2
        left_text = '\n'.join(all_paragraphs[:mid_point])
        right_text = '\n'.join(all_paragraphs[mid_point:])
        
        # Apply bullet points formatting for both sides
        set_bullet_points(left_content.text_frame, left_text, 37, 20)
        set_bullet_points(right_content.text_frame, right_text, 37, 20)
        
        return slide
    else:
        # Content will overflow, put max_per_side on each side
        left_text = '\n'.join(all_paragraphs[:max_per_side])
        right_text = '\n'.join(all_paragraphs[max_per_side:max_per_side*2])
        
        # Apply bullet points formatting for both sides
        set_bullet_points(left_content.text_frame, left_text, 37, 20)
        set_bullet_points(right_content.text_frame, right_text, 37, 20)
        
        # Create continuation slide with remaining content
        if len(all_paragraphs) > max_per_side * 2:
            continuation_data = slide_data.copy()
            continuation_data["content"] = '\n'.join(all_paragraphs[max_per_side*2:])
            add_two_content_slide(prs, continuation_data, True)
        
        return slide

def add_section_header_slide(prs, slide_data):
    """Add a section header slide."""
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="section_header")

    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slide_data["title"]
    
    # Process newlines in the subheading but maintain paragraphs
    subtitle.text_frame.clear()
    
    # Check which key exists in the slide_data
    sub_heading_key = "sub_heading" if "sub_heading" in slide_data else "sub-heading"
    
    # Split by explicit newlines in the input
    sub_heading_paragraphs = slide_data.get(sub_heading_key, "").split('\n')
    
    # Check for overflow in subheading
    first_part, overflow = check_content_overflow(sub_heading_paragraphs, MAX_BULLET_POINTS_PER_SLIDE // 2)
    
    for i, paragraph in enumerate(first_part):
        para = subtitle.text_frame.add_paragraph() if i > 0 else subtitle.text_frame.paragraphs[0]
        para.text = format_text(paragraph, 90)
        para.font.size = Pt(22)
    
    # Handle overflow if present
    if overflow:
        # Create continuation section header
        overflow_slide_data = slide_data.copy()
        overflow_slide_data["title"] = f"{slide_data['title']} (cont.)"
        if "sub_heading" in slide_data:
            overflow_slide_data["sub_heading"] = '\n'.join(overflow)
        else:
            overflow_slide_data["sub-heading"] = '\n'.join(overflow)
        add_section_header_slide(prs, overflow_slide_data)
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(40)
    title_format.bold = True
    
    return slide

def add_comparison_slide(prs, slide_data, is_continuation=False):
    """Add a comparison slide with two columns and overflow handling."""
    slide_layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="comparison")

    
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
        
    # Get the left content and handle potential key differences
    left_side_key = next((k for k in ["left_content", "left_layout"] if k in slide_data), None)
    right_side_key = next((k for k in ["right_content", "right_layout"] if k in slide_data), None)
    
    if not left_side_key or not right_side_key:
        # Fallback for invalid data
        left_content.text_frame.text = "Missing left content data"
        right_content.text_frame.text = "Missing right content data"
        title_format = title.text_frame.paragraphs[0].font
        title_format.size = Pt(36)
        title_format.bold = True
        return slide
    
    # Process left side
    left_title = slide_data[left_side_key].get("title", "Left Side")
    left_text = slide_data[left_side_key].get("content", "")
    
    # Process right side
    right_title = slide_data[right_side_key].get("title", "Right Side")
    right_text = slide_data[right_side_key].get("content", "")
    
    # Split content by newlines for overflow detection
    left_paragraphs = left_text.split('\n') if left_text else []
    right_paragraphs = right_text.split('\n') if right_text else []
    
    # Determine max bullets per column
    max_per_column = MAX_BULLET_POINTS_PER_SLIDE // 2 - 1  # Subtract 1 for title
    
    # Check for overflow
    left_first, left_overflow = check_content_overflow(left_paragraphs, max_per_column)
    right_first, right_overflow = check_content_overflow(right_paragraphs, max_per_column)
    
    # Add content with titles to left side
    left_content.text_frame.clear()
    title_para = left_content.text_frame.paragraphs[0]
    title_para.text = left_title
    title_para.font.bold = True
    
    # Add content with bullet points respecting newlines
    for paragraph in left_first:
        para = left_content.text_frame.add_paragraph()
        para.text = format_text(paragraph, 30)
        para.level = 0
        para.font.size = Pt(20)
    
    # Same for right side
    right_content.text_frame.clear()
    title_para = right_content.text_frame.paragraphs[0]
    title_para.text = right_title
    title_para.font.bold = True
    
    # Add content with bullet points respecting newlines
    for paragraph in right_first:
        para = right_content.text_frame.add_paragraph()
        para.text = format_text(paragraph, 30)
        para.level = 0
        para.font.size = Pt(20)
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(36)
    title_format.bold = True
    
    # Handle overflow if present
    if left_overflow or right_overflow:
        # Create continuation slide
        continuation_data = slide_data.copy()
        
        # Update content for continuation
        continuation_data[left_side_key] = continuation_data[left_side_key].copy()
        continuation_data[right_side_key] = continuation_data[right_side_key].copy()
        
        if left_overflow:
            continuation_data[left_side_key]["content"] = '\n'.join(left_overflow)
        else:
            continuation_data[left_side_key]["content"] = ""
            
        if right_overflow:
            continuation_data[right_side_key]["content"] = '\n'.join(right_overflow)
        else:
            continuation_data[right_side_key]["content"] = ""
            
        add_comparison_slide(prs, continuation_data, True)
    
    return slide

# def add_content_with_caption_slide(prs, slide_data, is_continuation=False):
#     """Add a slide with bullet point content, image, and caption — with overflow handling."""
#     slide_layout = prs.slide_layouts[1]  # Title and Content
#     slide = prs.slides.add_slide(slide_layout)
#     add_layout_styled_background(prs, slide, layout_type="content_with_caption")

#     title_shape = slide.shapes.title
#     content_shape = slide.placeholders[1]
    
#     content_data = slide_data.get("content", {})
#     if isinstance(content_data, str):
#         content_data = {"title": slide_data.get("title", ""), "content": content_data}

#     # Set title text (adjusted for continuation)
#     main_title = content_data.get("title", slide_data.get("title", ""))
#     title_shape.text = f"{main_title} (cont.)" if is_continuation else main_title
#     title_shape.text_frame.paragraphs[0].font.size = Pt(36)
#     title_shape.text_frame.paragraphs[0].font.bold = True

#     # Prepare content and handle overflow
#     content_text = content_data.get("content", "")
#     paragraphs = content_text.split('\n')
#     first_part, overflow = check_content_overflow(paragraphs)

#     # Add bullet points
#     # set_bullet_points(content_shape.text_frame, '\n'.join(first_part), font_size=Pt(20), left_indent=Inches(0.5))
#     set_bullet_points(content_shape.text_frame, '\n'.join(first_part), font_size=(20))


#     if not is_continuation:
#         # Image handling
#         image_path = slide_data.get("multi_media", slide_data.get("image_path", None))
#         left = Inches(5.8)
#         top = Inches(2)
#         width = Inches(3)
#         height = Inches(2.5)

#         if image_path:
#             try:
#                 slide.shapes.add_picture(image_path, left, top, width, height)
#             except Exception:
#                 # Add image path as placeholder text if invalid
#                 placeholder_box = slide.shapes.add_textbox(left, top, width, height)
#                 placeholder_frame = placeholder_box.text_frame
#                 placeholder_frame.text = f"[Image: {image_path}]"
#                 placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
#                 placeholder_frame.paragraphs[0].font.italic = True
#                 placeholder_frame.paragraphs[0].font.size = Pt(14)

#         # Caption (if available)
#         caption = content_data.get("caption", "")
#         if caption:
#             caption_top = top + height + Inches(0.2)
#             caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.5))
#             caption_frame = caption_box.text_frame
#             caption_frame.text = caption
#             caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
#             caption_frame.paragraphs[0].font.italic = True
#             caption_frame.paragraphs[0].font.size = Pt(14)

#     # Recursive call for overflow
#     if overflow:
#         continuation_data = slide_data.copy()
#         if isinstance(continuation_data.get("content", {}), dict):
#             continuation_data["content"] = continuation_data["content"].copy()
#             continuation_data["content"]["content"] = '\n'.join(overflow)
#         else:
#             continuation_data["content"] = {'title': main_title, 'content': '\n'.join(overflow)}
#         add_content_with_caption_slide(prs, continuation_data, is_continuation=True)

#     return slide

def add_content_with_caption_slide(prs, slide_data, is_continuation=False):
    """Add a slide with bullet point content and optional icon/3D image."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="content_with_caption")

    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    content_data = slide_data.get("content", {})
    if isinstance(content_data, str):
        content_data = {"title": slide_data.get("title", ""), "content": content_data}

    # Set title with appropriate continuation marker
    main_title = content_data.get("title", slide_data.get("title", ""))
    title_shape.text = f"{main_title} (cont.)" if is_continuation else main_title
    
    # Apply consistent formatting to all title paragraphs
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
    
    # Process content with line breaks at 45 characters
    content_text = content_data.get("content", "")
    formatted_content = []
    for line in content_text.split('\n'):
        words = line.split()
        current_line = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + (1 if current_length > 0 else 0) <= 45:
                current_line.append(word)
                current_length += len(word) + (1 if current_length > 0 else 0)
            else:
                formatted_content.append(" ".join(current_line))
                current_line = [word]
                current_length = len(word)
        
        if current_line:
            formatted_content.append(" ".join(current_line))
    
    # Check for overflow using MAX_BULLET_POINTS_PER_SLIDE
    first_part, overflow = check_content_overflow(formatted_content)
    
    # Add bullet points using the referenced function
    set_bullet_points(content_shape.text_frame, '\n'.join(first_part), max_line_length=45, font_size=20)

    if not is_continuation:
        # 3D icon or chart handling
        chart_data = slide_data.get("chart/smart3D_icon")
        if chart_data:
            left = Inches(5.8)
            top = Inches(2)
            width = Inches(3)
            height = Inches(2.5)
            
            # Add placeholder or actual chart/3D based on implementation
            icon_box = slide.shapes.add_textbox(left, top, width, height)
            icon_frame = icon_box.text_frame
            icon_frame.text = f"[3D Icon/Chart]"
            para = icon_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(14)

    # Recursive call for overflow
    if overflow:
        continuation_data = slide_data.copy()
        if isinstance(continuation_data.get("content", {}), dict):
            continuation_data["content"] = continuation_data["content"].copy()
            continuation_data["content"]["content"] = '\n'.join(overflow)
        else:
            continuation_data["content"] = {'title': main_title, 'content': '\n'.join(overflow)}
        add_content_with_caption_slide(prs, continuation_data, is_continuation=True)

    return slide

def add_image_with_caption_slide(prs, slide_data):
    """Add a slide with an image and caption."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="image_with_caption")

    # Add title with max 60 characters, line breaks at 30 chars
    title_text = slide_data.get("title", "")
    
    # Format title with line breaks
    formatted_title_lines = []
    words = title_text.split()
    current_line = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + (1 if current_length > 0 else 0) <= 30:
            current_line.append(word)
            current_length += len(word) + (1 if current_length > 0 else 0)
        else:
            formatted_title_lines.append(" ".join(current_line))
            current_line = [word]
            current_length = len(word)
    
    if current_line:
        formatted_title_lines.append(" ".join(current_line))
    
    formatted_title = "\n".join(formatted_title_lines)
    
    # Add title textbox
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = formatted_title
    
    # Apply consistent formatting to all title paragraphs
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    # Add image
    image_path = slide_data.get("image_path", "")
    left = Inches(2)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(4)
    
    try:
        slide.shapes.add_picture(image_path, left, top, width, height)
    except Exception:
        # Add image path as placeholder text if invalid
        placeholder_box = slide.shapes.add_textbox(left, top, width, height)
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.text = f"[Image: {image_path}]"
        paragraph = placeholder_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.italic = True
        paragraph.font.size = Pt(14)

    # Add caption with max 250 characters
    caption_text = slide_data.get("content", "")
    if len(caption_text) > 250:
        caption_text = caption_text[:247] + "..."

    # Format caption with line breaks at 90 characters
    formatted_caption_lines = []
    words = caption_text.split()
    current_line = []
    current_length = 0

    for word in words:
        if current_length + len(word) + (1 if current_length > 0 else 0) <= 90:
            current_line.append(word)
            current_length += len(word) + (1 if current_length > 0 else 0)
        else:
            formatted_caption_lines.append(" ".join(current_line))
            current_line = [word]
            current_length = len(word)

    if current_line:
        formatted_caption_lines.append(" ".join(current_line))

    formatted_caption = "\n".join(formatted_caption_lines)

    caption_top = top + height + Inches(0.3)
    caption_box = slide.shapes.add_textbox(Inches(1), caption_top, Inches(8), Inches(1))
    caption_frame = caption_box.text_frame
    caption_frame.text = formatted_caption

    # Format all caption paragraphs
    for paragraph in caption_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.italic = True
        paragraph.font.size = Pt(16)

    return slide

def add_table_slide(prs, slide_data):
    """Add a slide with a table using the default title and dynamic sizing."""
    slide_layout = prs.slide_layouts[5]  # Using a blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background styling (optional)
    add_layout_styled_background(prs, slide, layout_type="title_with_table")

    # Use the title placeholder if available
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Table Slide")
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.LEFT

    # Get table data
    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers or not rows:
        return slide

    n_rows = len(rows) + 1
    n_cols = len(headers)

    # Set table position
    left = Inches(0.5)
    top = Inches(1.5)
    max_width = Inches(9)
    max_height = Inches(5)

    # Estimate column widths based on max string length
    max_lengths = [len(str(header)) for header in headers]
    for row in rows:
        for i, cell in enumerate(row):
            max_lengths[i] = max(max_lengths[i], len(str(cell)))

    total_length = sum(max_lengths)
    col_widths = [max_width * (l / total_length) for l in max_lengths]

    # Add table with approximate dimensions
    table = slide.shapes.add_table(n_rows, n_cols, left, top, max_width, max_height).table

    # Apply calculated widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = int(width)

    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(14)
            para.alignment = PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.alignment = PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Dynamic row heights based on number of lines
    for i in range(n_rows):
        row_height = Inches(0.4)
        for j in range(n_cols):
            lines = table.cell(i, j).text.count('\n') + 1
            row_height = max(row_height, Inches(0.2 + 0.2 * lines))
        table.rows[i].height = int(row_height)

    # Optional caption/content below the table
    content_text = slide_data.get("content", "")
    if content_text:
        content_top = top + max_height + Inches(0.3)
        content_box = slide.shapes.add_textbox(left, content_top, max_width, Inches(2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.clear()  # Clear default paragraph

        for line in content_text.strip().split('\n'):
            p = content_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT

    return slide

def format_text(text, max_line_length):
    """Format text to wrap at specified character length, keeping it as a single paragraph."""
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + len(current_line) <= max_line_length:
            current_line.append(word)
            current_length += len(word)
        else:
            lines.append(' '.join(current_line))
            current_line = [word]
            current_length = len(word)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    # Join the lines with spaces rather than newlines to keep in same paragraph
    return ' '.join(lines)

def create_presentation_from_json(json_data):
    """Create a presentation from JSON data."""
    prs = Presentation()
    
    for slide_data in json_data:
        layout = slide_data.get("layout", "")
        
        if layout == "title_only":
            add_title_only_slide(prs, slide_data)
        elif layout == "title_slide":
            add_title_slide(prs, slide_data)
        elif layout == "title_and_content":
            add_title_and_content_slide(prs, slide_data)
        elif layout == "two_content":
            add_two_content_slide(prs, slide_data)
        elif layout == "section_header":
            add_section_header_slide(prs, slide_data)
        elif layout == "comparison":
            add_comparison_slide(prs, slide_data)
        elif layout in "content_with_caption":
            add_content_with_caption_slide(prs, slide_data)
        elif layout in "image_with_caption":
            add_image_with_caption_slide(prs, slide_data)
        elif layout in ["title_with_table", "chart", "other_multi_media"]:
            add_table_slide(prs, slide_data)
    
    return prs

def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint from JSON')
    parser.add_argument('json_file', help='Path to the JSON file')
    parser.add_argument('--output', default='presentation.pptx', help='Output file name')
    parser.add_argument('--max-bullets', type=int, default=MAX_BULLET_POINTS_PER_SLIDE, 
                        help=f'Maximum number of bullet points per slide (default: {MAX_BULLET_POINTS_PER_SLIDE})')
    
    args = parser.parse_args()
    
    # Update global max bullets if specified
    # global MAX_BULLET_POINTS_PER_SLIDE
    # if args.max_bullets:
    #     MAX_BULLET_POINTS_PER_SLIDE = args.max_bullets
    
    try:
        with open(args.json_file, 'r', encoding='utf-8') as file:
            json_data = json.load(file)
        
        prs = create_presentation_from_json(json_data)
        prs.save(args.output)
        
        print(f"Presentation created successfully: {args.output}")
        
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    main()