# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# import os
# import json

# # ========== SHARED UTILS ==========
# def add_gradient_overlay(slide, color1=RGBColor(255, 255, 255), color2=RGBColor(240, 240, 255)):
#     for transparency, color in [(0.0, color1), (0.3, color2)]:
#         shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
#         fill = shape.fill
#         fill.solid()
#         fill.fore_color.rgb = color
#         fill.transparency = transparency
#         shape.line.fill.background()

# def estimate_lines(text, max_chars_per_line=45):
#     return max(1, len(text) // max_chars_per_line + (1 if len(text) % max_chars_per_line > 0 else 0))

# # ========== TITLE SLIDE ==========
# def create_dynamic_title_slide(prs, title, subtitle):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     add_gradient_overlay(slide)

#     heading_color = RGBColor(20, 33, 61)
#     subheading_color = RGBColor(60, 60, 60)

#     left_margin = Inches(1.0)
#     width = Inches(8.0)

#     title_lines = estimate_lines(title)
#     title_height = Inches(1.0 + 0.5 * title_lines)
#     top_title = Inches(1.5 - (0.2 * (title_lines - 1))) if title_lines == 1 else Inches(1.0)

#     title_box = slide.shapes.add_textbox(left_margin, top_title, width, title_height)
#     tf_title = title_box.text_frame
#     tf_title.word_wrap = True
#     tf_title.clear()
#     p_title = tf_title.paragraphs[0]
#     run_title = p_title.add_run()
#     run_title.text = title
#     p_title.alignment = PP_ALIGN.CENTER
#     run_title.font.name = 'Arial'
#     run_title.font.size = Pt(40 if title_lines == 1 else 34)
#     run_title.font.bold = True
#     run_title.font.color.rgb = heading_color

#     subtitle_top = top_title + title_height + Inches(0.2)
#     subtitle_box = slide.shapes.add_textbox(left_margin + Inches(0.2), subtitle_top, width - Inches(0.4), Inches(1))
#     tf_sub = subtitle_box.text_frame
#     tf_sub.word_wrap = True
#     tf_sub.clear()
#     p_sub = tf_sub.paragraphs[0]
#     run_sub = p_sub.add_run()
#     run_sub.text = subtitle
#     p_sub.alignment = PP_ALIGN.CENTER
#     run_sub.font.name = 'Arial'
#     run_sub.font.size = Pt(24)
#     run_sub.font.color.rgb = subheading_color

# # ========== TEXT + IMAGE SLIDE ==========
# def create_text_and_image_layout(slide, title, subheading, bullet_points, image_path, bullet_icon="➤"):
#     heading_color = RGBColor(20, 33, 61)
#     subheading_color = RGBColor(60, 60, 60)
#     bullet_color = RGBColor(40, 40, 40)

#     title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
#     text_frame = title_box.text_frame
#     text_frame.clear()
#     p = text_frame.paragraphs[0]
#     run = p.add_run()
#     run.text = title
#     p.alignment = PP_ALIGN.CENTER
#     run.font.name = 'Arial'
#     run.font.size = Pt(38)
#     run.font.bold = True
#     run.font.color.rgb = heading_color

#     content_top = Inches(1.4)
#     content_height = Inches(5.6)
#     text_left = Inches(0.7)
#     text_width = Inches(4.8)
#     image_left = Inches(5.7)
#     image_width = Inches(3.5)

#     image_found = image_path and os.path.isfile(image_path)

#     if image_found:
#         try:
#             slide.shapes.add_picture(image_path, image_left, content_top, height=content_height)
#         except Exception as e:
#             print(f"⚠️ Error loading image: {e}")
#             image_found = False

#     text_box = slide.shapes.add_textbox(text_left, content_top, text_width, content_height)
#     tf = text_box.text_frame
#     tf.word_wrap = True
#     tf.margin_top = 0
#     tf.margin_bottom = 0

#     p_sub = tf.paragraphs[0]
#     run_sub = p_sub.add_run()
#     run_sub.text = subheading
#     p_sub.alignment = PP_ALIGN.LEFT
#     run_sub.font.name = 'Arial'
#     run_sub.font.size = Pt(22)
#     run_sub.font.bold = True
#     run_sub.font.color.rgb = subheading_color
#     p_sub.space_after = Pt(12)

#     for point in bullet_points:
#         p = tf.add_paragraph()
#         p.text = f"{bullet_icon} {point}"
#         p.level = 0
#         p.font.size = Pt(20)
#         p.font.name = 'Arial'
#         p.font.color.rgb = bullet_color
#         p.alignment = PP_ALIGN.LEFT
#         p.space_after = Pt(6)

# def create_styled_slide(prs, title, subheading, bullet_points, image_path, bullet_icon="➤"):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     add_gradient_overlay(slide)
#     create_text_and_image_layout(slide, title, subheading, bullet_points, image_path, bullet_icon)

# # ========== CLEAN SLIDE (NO IMAGE) ==========
# def create_clean_slide(prs, title, bullet_points, bullet_icon="➤"):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     add_gradient_overlay(slide)

#     heading_color = RGBColor(20, 33, 61)
#     bullet_color = RGBColor(40, 40, 40)

#     title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
#     tf = title_box.text_frame
#     tf.clear()
#     p = tf.paragraphs[0]
#     run = p.add_run()
#     run.text = title
#     p.alignment = PP_ALIGN.CENTER
#     run.font.name = 'Arial'
#     run.font.size = Pt(36)
#     run.font.bold = True
#     run.font.color.rgb = heading_color

#     content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8), Inches(5.5))
#     tf_content = content_box.text_frame
#     tf_content.word_wrap = True

#     for point in bullet_points:
#         p = tf_content.add_paragraph()
#         p.text = f"{bullet_icon} {point}"
#         p.level = 0
#         p.font.size = Pt(22)
#         p.font.name = 'Arial'
#         p.font.color.rgb = bullet_color
#         p.alignment = PP_ALIGN.LEFT
#         p.space_after = Pt(8)

# # ========== MAIN PPT BUILDER ==========
# def build_combined_ppt_from_json(json_path, output_path="auto_selected_layout.pptx"):
#     prs = Presentation()

#     with open(json_path, 'r', encoding='utf-8') as f:
#         slide_data_list = json.load(f)

#     for slide_data in slide_data_list:
#         layout_type = slide_data.get("layout")
#         title = slide_data.get("title", "")
#         content = slide_data.get("content", [])
#         image_path = slide_data.get("image_path", "")
#         bullet_icon = slide_data.get("icon", "•")

#         if layout_type == "Title Slide":
#             subtitle = "\n".join(content)
#             create_dynamic_title_slide(prs, title, subtitle)

#         elif layout_type == "Title and Content":
#             create_clean_slide(
#                 prs,
#                 title=title,
#                 bullet_points=content,
#                 bullet_icon=bullet_icon
#             )
#         elif layout_type == "Text + Image":
#             subheading = slide_data.get("subheading", "")
#             create_styled_slide(
#                 prs,
#                 title=title,
#                 subheading=subheading,
#                 bullet_points=content,
#                 image_path=image_path,
#                 bullet_icon=bullet_icon
#             )
#         else:
#             print(f"⚠️ Unknown layout: {layout_type}, skipping...")

#     prs.save(output_path)
#     print(f"✅ Presentation saved locally as: {output_path}")

# # ========== RUN ==========

# if __name__ == "__main__":
#     # Update these paths
#     input_json = "slides.json"
#     output_pptx = "final_presentation.pptx"
#     build_combined_ppt_from_json(input_json, output_pptx)





from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json
try:
    from PIL import Image
except ImportError:
    print("⚠️ PIL library not found. Install with: pip install pillow")
    print("    Image aspect ratio handling will be disabled.")

# ========== CONTENT OVERFLOW HANDLER ==========
def handle_content_overflow(slide, content, container_width, container_height, content_type="text", **kwargs):
    """
    Comprehensive function to handle content overflow in PowerPoint slides.
    
    Args:
        slide: The slide object where content is being added
        content: The content to add (text string or list for text, image path for images)
        container_width: Available width in inches
        container_height: Available height in inches
        content_type: "text" or "image"
        **kwargs: Additional parameters based on content_type
    """
    # Default parameters
    defaults = {
        'text': {
            'font_name': 'Arial',
            'base_font_size': 20,  
            'min_font_size': 12,
            'bullet_icon': '•',
            'alignment': PP_ALIGN.LEFT,
            'is_title': False,
            'bold': False,
            'left_position': 0,
            'top_position': 0,
            'color': RGBColor(0, 0, 0)
        },
        'image': {
            'maintain_aspect': True,
            'left_position': 0,
            'top_position': 0
        }
    }
    
    # Merge defaults with provided kwargs
    params = defaults[content_type].copy()
    for k, v in kwargs.items():
        if k in params:
            params[k] = v
    
    # Handle text content
    if content_type == "text":
        # Create text box
        text_box = slide.shapes.add_textbox(
            Inches(params['left_position']), 
            Inches(params['top_position']), 
            Inches(container_width), 
            Inches(container_height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        # Process content based on type
        if isinstance(content, list):
            # Handle bullet points
            all_lines = []
            
            # First determine if we need to split lines
            for point in content:
                all_lines.extend(split_long_line(point, max_length=int(container_width * 12)))
            
            # Calculate appropriate font size based on number of lines
            max_lines = int(container_height / 0.3)  # Rough estimate: 0.3 inches per line
            font_size = get_font_size_for_lines(
                all_lines, 
                max_lines=max_lines,
                base_size=params['base_font_size'], 
                min_size=params['min_font_size']
            )
            
            # Add the paragraphs with calculated font size
            first_p = True
            for line in all_lines:
                if first_p:
                    p = tf.paragraphs[0]
                    first_p = False
                else:
                    p = tf.add_paragraph()
                
                p.text = f"{params['bullet_icon']} {line}" if not params['is_title'] else line
                p.font.name = params['font_name']
                p.font.size = Pt(font_size)
                p.font.bold = params['bold']
                p.font.color.rgb = params['color']
                p.alignment = params['alignment']
                p.space_after = Pt(4)
                
        else:
            # Handle single text block
            lines = split_long_line(content, max_length=int(container_width * 12))
            max_lines = int(container_height / 0.3)
            font_size = get_font_size_for_lines(
                lines, 
                max_lines=max_lines,
                base_size=params['base_font_size'], 
                min_size=params['min_font_size']
            )
            
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = content
            run.font.name = params['font_name']
            run.font.size = Pt(font_size)
            run.font.bold = params['bold']
            run.font.color.rgb = params['color']
            p.alignment = params['alignment']
        
        return text_box
    
    # Handle image content
    elif content_type == "image":
        if not content or not os.path.isfile(content):
            print(f"⚠️ Image not found: {content}")
            return None
            
        try:
            # Calculate dimensions that maintain aspect ratio and fit container
            if params['maintain_aspect']:
                try:
                    # Get image dimensions using PIL
                    with Image.open(content) as img:
                        img_width, img_height = img.size
                    
                    # Calculate aspect ratio
                    aspect_ratio = img_width / img_height
                    
                    # Determine if width or height is the limiting factor
                    container_aspect = container_width / container_height
                    
                    if aspect_ratio > container_aspect:
                        # Width limited
                        width = container_width
                        height = width / aspect_ratio
                    else:
                        # Height limited
                        height = container_height
                        width = height * aspect_ratio
                        
                    # Add picture with calculated dimensions
                    picture = slide.shapes.add_picture(
                        content, 
                        Inches(params['left_position']), 
                        Inches(params['top_position']), 
                        width=Inches(width), 
                        height=Inches(height)
                    )
                except (NameError, ImportError):
                    # PIL not available, fall back to height-based scaling
                    picture = slide.shapes.add_picture(
                        content, 
                        Inches(params['left_position']), 
                        Inches(params['top_position']), 
                        height=Inches(container_height)
                    )
            else:
                # Add picture with full container dimensions
                picture = slide.shapes.add_picture(
                    content, 
                    Inches(params['left_position']), 
                    Inches(params['top_position']), 
                    width=Inches(container_width), 
                    height=Inches(container_height)
                )
            
            return picture
            
        except Exception as e:
            print(f"⚠️ Image processing error: {e}")
            return None
    
    else:
        print(f"⚠️ Unknown content type: {content_type}")
        return None

# ========== HELPER FUNCTIONS ==========
def split_long_line(text, max_length=90):
    """Split text into multiple lines if it exceeds max_length"""
    if not isinstance(text, str):
        text = str(text)
        
    words = text.split()
    lines = []
    current_line = ""
    
    for word in words:
        # If adding this word would exceed max_length, start a new line
        if len(current_line) + len(word) + 1 <= max_length:
            current_line += (" " if current_line else "") + word
        else:
            lines.append(current_line)
            current_line = word
            
    # Add the last line if not empty
    if current_line:
        lines.append(current_line)
        
    return lines

def get_font_size_for_lines(lines, max_lines, base_size=20, min_size=12):
    """Calculate appropriate font size based on number of lines"""
    if not lines:
        return base_size
        
    # Scale down font size if there are too many lines
    scale = min(1.0, max_lines / len(lines))
    return max(min_size, int(base_size * scale))

def get_font_size_for_text(text, max_chars, base_size=40, min_size=12):
    """Calculate appropriate font size based on text length"""
    scale = min(1.0, max_chars / max(len(text), 1))
    return max(min_size, int(base_size * scale))

def add_gradient_overlay(slide, color1=RGBColor(255, 255, 255), color2=RGBColor(240, 240, 255)):
    """Add a gradient background overlay to slide"""
    for transparency, color in [(0.0, color1), (0.3, color2)]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
        fill.transparency = transparency
        shape.line.fill.background()

# ========== SLIDE CREATION FUNCTIONS ==========
def create_dynamic_title_slide(prs, title, subtitle):
    """Create a title slide with proper overflow handling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_overlay(slide)

    heading_color = RGBColor(20, 33, 61)
    subheading_color = RGBColor(60, 60, 60)

    # Add title with overflow handling
    handle_content_overflow(
        slide, 
        title, 
        container_width=8.0, 
        container_height=1.2, 
        content_type="text",
        left_position=1.0,
        top_position=1.0,
        base_font_size=40,
        min_font_size=28,
        is_title=True,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        color=heading_color
    )

    # Add subtitle with overflow handling
    handle_content_overflow(
        slide, 
        subtitle, 
        container_width=7.6, 
        container_height=1.0, 
        content_type="text",
        left_position=1.2,
        top_position=2.2,
        base_font_size=24,
        min_font_size=16,
        is_title=False,
        bold=False,
        alignment=PP_ALIGN.CENTER,
        color=subheading_color
    )

def create_text_and_image_layout(slide, title, subheading, bullet_points, image_path, bullet_icon="➤"):
    """Create a text and image slide with proper overflow handling"""
    heading_color = RGBColor(20, 33, 61)
    subheading_color = RGBColor(60, 60, 60)
    bullet_color = RGBColor(40, 40, 40)

    # Add title with overflow handling
    handle_content_overflow(
        slide, 
        title, 
        container_width=9.0, 
        container_height=0.9, 
        content_type="text",
        left_position=0.5,
        top_position=0.3,
        base_font_size=34,
        min_font_size=20,
        is_title=True,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        color=heading_color
    )

    # Add image with overflow handling if path exists
    if image_path:
        handle_content_overflow(
            slide,
            image_path, 
            container_width=3.5, 
            container_height=5.6, 
            content_type="image",
            left_position=5.7,
            top_position=1.4,
            maintain_aspect=True
        )

    # Add subheading
    subheading_box = handle_content_overflow(
        slide, 
        subheading, 
        container_width=4.8, 
        container_height=0.7, 
        content_type="text",
        left_position=0.7,
        top_position=1.4,
        base_font_size=22,
        min_font_size=14,
        is_title=False,
        bold=True,
        alignment=PP_ALIGN.LEFT,
        color=subheading_color
    )

    # Add bullet points
    handle_content_overflow(
        slide,
        bullet_points, 
        container_width=4.8, 
        container_height=4.9, 
        content_type="text",
        left_position=0.7,
        top_position=2.1,
        base_font_size=20,
        min_font_size=12,
        bullet_icon=bullet_icon,
        color=bullet_color
    )

def create_clean_slide(prs, title, bullet_points, bullet_icon="➤"):
    """Create a clean slide with text only and proper overflow handling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_overlay(slide)

    heading_color = RGBColor(20, 33, 61)
    bullet_color = RGBColor(40, 40, 40)

    # Add title with overflow handling
    handle_content_overflow(
        slide, 
        title, 
        container_width=9.0, 
        container_height=0.9, 
        content_type="text",
        left_position=0.5,
        top_position=0.4,
        base_font_size=34,
        min_font_size=24,
        is_title=True,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        color=heading_color
    )

    # Add bullet points with overflow handling
    handle_content_overflow(
        slide,
        bullet_points, 
        container_width=8.0, 
        container_height=5.3, 
        content_type="text",
        left_position=1.0,
        top_position=1.5,
        base_font_size=20,
        min_font_size=12,
        bullet_icon=bullet_icon,
        color=bullet_color
    )

def create_styled_slide(prs, title, subheading, bullet_points, image_path, bullet_icon="➤"):
    """Create a styled slide with text and image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_overlay(slide)
    create_text_and_image_layout(slide, title, subheading, bullet_points, image_path, bullet_icon)

# ========== MAIN PPT BUILDER ==========
def build_combined_ppt_from_json(json_path, output_path="auto_selected_layout.pptx"):
    """Build a PowerPoint presentation from JSON data"""
    prs = Presentation()

    try:
        # Open file using utf-8 encoding
        with open(json_path, 'r', encoding='utf-8') as f:
            slide_data_list = json.load(f)
            
        if not slide_data_list:
            print("⚠️ Error: JSON file contains no slide data.")
            return
    except Exception as e:
        print(f"⚠️ Error loading JSON file: {e}")
        return

    for slide_idx, slide_data in enumerate(slide_data_list):
        try:
            layout_type = slide_data.get("layout", "")
            title = slide_data.get("title", "")
            content = slide_data.get("content", [])
            image_path = slide_data.get("image_path", "")
            bullet_icon = slide_data.get("icon", "•")
            
            print(f"Processing slide {slide_idx+1}: {layout_type}")

            if layout_type == "Title Slide":
                subtitle = "\n".join(content) if isinstance(content, list) else str(content)
                create_dynamic_title_slide(prs, title, subtitle)

            elif layout_type == "Title and Content":
                create_clean_slide(
                    prs,
                    title=title,
                    bullet_points=content if isinstance(content, list) else [str(content)],
                    bullet_icon=bullet_icon
                )
                
            elif layout_type == "Text + Image":
                subheading = slide_data.get("subheading", "")
                create_styled_slide(
                    prs,
                    title=title,
                    subheading=subheading,
                    bullet_points=content if isinstance(content, list) else [str(content)],
                    image_path=image_path,
                    bullet_icon=bullet_icon
                )
                
            else:
                print(f"⚠️ Unknown layout type: '{layout_type}' on slide {slide_idx+1}, skipping.")
                
        except Exception as e:
            print(f"⚠️ Error processing slide {slide_idx+1}: {e}")
            # Continue with other slides instead of crashing

    try:
        prs.save(output_path)
        print(f"✅ Presentation saved successfully as: {output_path}")
    except Exception as e:
        print(f"⚠️ Error saving presentation: {e}")

# ========== EXECUTE IF MAIN ==========
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        json_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else "presentation.pptx"
        build_combined_ppt_from_json(json_path, output_path)
    else:
        # Default paths if no arguments provided
        json_path = input("Enter path to JSON file: ")
        output_path = input("Enter output PPTX path (or press Enter for default): ") or "presentation.pptx"
        build_combined_ppt_from_json(json_path, output_path)