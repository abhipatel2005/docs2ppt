from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from style import add_layout_styled_background

# Configuration constants
MAX_BULLET_POINTS_PER_SLIDE = 8  # Maximum number of bullet points per slide
MAX_CONTENT_HEIGHT = 5  # Maximum content height in inches

#helper function
def set_bullet_points(text_frame, content, max_line_length=80, font_size=20):
    text_frame.clear()
    
    # First split by explicit newlines in the content
    paragraphs = content.split('\n')
    
    for i, paragraph in enumerate(paragraphs):
        # Add a new paragraph for each explicit newline in the input
        para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        
        # Format text to respect line length but keep all text in the same paragraph
        formatted_text = format_text(paragraph, max_line_length)
        para.text = formatted_text
        para.level = 0
        para.font.size = Pt(font_size)
    
    return len(paragraphs)  # Return the number of bullet points added

def check_content_overflow(paragraphs, max_bullets=MAX_BULLET_POINTS_PER_SLIDE):
    """Check if content will overflow the slide and split if needed."""
    if len(paragraphs) <= max_bullets:
        return paragraphs, None
    
    # Split content
    first_part = paragraphs[:max_bullets]
    overflow_part = paragraphs[max_bullets:]
    
    return first_part, overflow_part

def add_title_only_slide(prs, slide_data):
    """Add a slide with only a title."""
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_only")
    
    title = slide.shapes.title
    title.text = slide_data["title"]
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(44)
    title_format.bold = True
    
    return slide

def add_title_slide(prs, slide_data):
    """Add a title slide with title and subtitle."""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_slide")
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Set title text
    title.text = slide_data["title"]
    
    # Apply formatting to ALL paragraphs in the title text frame
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
    
    # Check which key exists in the slide_data
    sub_heading_key = "sub-heading" if "sub-heading" in slide_data else "sub_heading"
    subtitle.text = slide_data.get(sub_heading_key, "")
    
    # Apply formatting to ALL paragraphs in the subtitle text frame
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.italic = True
    
    return slide

def add_title_and_content_slide(prs, slide_data, is_continuation=False):
    """Add a slide with title and content with precise vertical overflow handling."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="title_and_content")
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
    
    # Format all paragraphs in the title text frame
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True

    # Calculate the available height in the content placeholder
    available_height = content.height
    
    # Split content by newlines
    content_paragraphs = slide_data["content"].split('\n')
    
    # Create a testing function to determine exactly how many paragraphs will fit
    def test_content_fit(paragraphs, font_size=20):
        # Create a temporary slide for testing
        temp_slide = prs.slides.add_slide(prs.slide_layouts[1])
        temp_content = temp_slide.placeholders[1]
        
        # Set the content
        test_content = '\n'.join(paragraphs)
        set_bullet_points(temp_content.text_frame, test_content, 80, font_size)
        
        # Check if content fits within available height
        text_height = 0
        for p in temp_content.text_frame.paragraphs:
            # Get line spacing safely, default to 1.0 if None
            line_spacing = 1.0  # Default spacing
            if hasattr(p, 'line_spacing') and p.line_spacing is not None:
                line_spacing = p.line_spacing
                
            # Convert font size to Pt if it's not already
            font_size_pt = font_size
            if not isinstance(font_size_pt, Pt):
                font_size_pt = Pt(font_size)
                
            # Calculate paragraph height safely
            para_height = font_size_pt.pt * line_spacing * 1.2  # Use .pt attribute to get the numeric value
            text_height += para_height
        
        # Remove the temporary slide
        prs.slides._sldIdLst.remove(temp_slide._element.sldId)
        
        return text_height <= available_height
    
    # Use a simpler approach to avoid errors with temporary slides
    # For most presentations, using a fixed number based on paragraph length works well
    
    # Calculate average paragraph length
    avg_chars = sum(len(p) for p in content_paragraphs) / max(1, len(content_paragraphs))
    
    # Determine paragraphs per slide based on length
    if avg_chars > 200:  # Very long paragraphs
        max_paragraphs = 3
    elif avg_chars > 100:  # Medium paragraphs
        max_paragraphs = 4
    elif avg_chars > 50:  # Short paragraphs
        max_paragraphs = 6
    else:  # Very short paragraphs
        max_paragraphs = 8
        
    # Be more conservative for continuation slides
    if is_continuation:
        max_paragraphs = max(1, max_paragraphs - 1)
        
    # Calculate how many paragraphs to use on this slide
    paragraphs_to_use = min(max_paragraphs, len(content_paragraphs))
    
    # Split content
    first_part = content_paragraphs[:paragraphs_to_use]
    overflow = content_paragraphs[paragraphs_to_use:] if paragraphs_to_use < len(content_paragraphs) else None
    
    # Add first part content
    first_content = '\n'.join(first_part)
    set_bullet_points(content.text_frame, first_content, 80, 20)
    
    # Handle overflow if present
    if overflow and len(overflow) > 0:
        # Create continuation slide with remaining content
        continuation_data = slide_data.copy()
        continuation_data["content"] = '\n'.join(overflow)
        add_title_and_content_slide(prs, continuation_data, True)
    
    return slide

def add_two_content_slide(prs, slide_data, is_continuation=False):
    """Add a slide with title and two content areas with overflow handling."""
    slide_layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="two_content")

    
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(36)
    title_format.bold = True
    
    # Split content by newlines
    all_paragraphs = slide_data["content"].split('\n')
    total_paragraphs = len(all_paragraphs)
    
    # For this layout, we'll display half the content on each side
    max_per_side = MAX_BULLET_POINTS_PER_SLIDE // 2
    
    # Check if total content will fit
    if total_paragraphs <= MAX_BULLET_POINTS_PER_SLIDE:
        # Content will fit, split between columns
        mid_point = total_paragraphs // 2
        left_text = '\n'.join(all_paragraphs[:mid_point])
        right_text = '\n'.join(all_paragraphs[mid_point:])
        
        # Apply bullet points formatting for both sides
        set_bullet_points(left_content.text_frame, left_text, 37, 20)
        set_bullet_points(right_content.text_frame, right_text, 37, 20)
        
        return slide
    else:
        # Content will overflow, put max_per_side on each side
        left_text = '\n'.join(all_paragraphs[:max_per_side])
        right_text = '\n'.join(all_paragraphs[max_per_side:max_per_side*2])
        
        # Apply bullet points formatting for both sides
        set_bullet_points(left_content.text_frame, left_text, 37, 20)
        set_bullet_points(right_content.text_frame, right_text, 37, 20)
        
        # Create continuation slide with remaining content
        if len(all_paragraphs) > max_per_side * 2:
            continuation_data = slide_data.copy()
            continuation_data["content"] = '\n'.join(all_paragraphs[max_per_side*2:])
            add_two_content_slide(prs, continuation_data, True)
        
        return slide

def add_section_header_slide(prs, slide_data):
    """Add a section header slide."""
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="section_header")

    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slide_data["title"]
    
    # Process newlines in the subheading but maintain paragraphs
    subtitle.text_frame.clear()
    
    # Check which key exists in the slide_data
    sub_heading_key = "sub_heading" if "sub_heading" in slide_data else "sub-heading"
    
    # Split by explicit newlines in the input
    sub_heading_paragraphs = slide_data.get(sub_heading_key, "").split('\n')
    
    # Check for overflow in subheading
    first_part, overflow = check_content_overflow(sub_heading_paragraphs, MAX_BULLET_POINTS_PER_SLIDE // 2)
    
    for i, paragraph in enumerate(first_part):
        para = subtitle.text_frame.add_paragraph() if i > 0 else subtitle.text_frame.paragraphs[0]
        para.text = format_text(paragraph, 90)
        para.font.size = Pt(22)
    
    # Handle overflow if present
    if overflow:
        # Create continuation section header
        overflow_slide_data = slide_data.copy()
        overflow_slide_data["title"] = f"{slide_data['title']} (cont.)"
        if "sub_heading" in slide_data:
            overflow_slide_data["sub_heading"] = '\n'.join(overflow)
        else:
            overflow_slide_data["sub-heading"] = '\n'.join(overflow)
        add_section_header_slide(prs, overflow_slide_data)
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(40)
    title_format.bold = True
    
    return slide

def add_comparison_slide(prs, slide_data, is_continuation=False):
    """Add a comparison slide with two columns and overflow handling."""
    slide_layout = prs.slide_layouts[3]  # Two Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background (prs, slide, layout_type="comparison")

    
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    # Adjust title for continuation slides
    if is_continuation:
        title.text = f"{slide_data['title']} (cont.)"
    else:
        title.text = slide_data["title"]
        
    # Get the left content and handle potential key differences
    left_side_key = next((k for k in ["left_content", "left_layout"] if k in slide_data), None)
    right_side_key = next((k for k in ["right_content", "right_layout"] if k in slide_data), None)
    
    if not left_side_key or not right_side_key:
        # Fallback for invalid data
        left_content.text_frame.text = "Missing left content data"
        right_content.text_frame.text = "Missing right content data"
        title_format = title.text_frame.paragraphs[0].font
        title_format.size = Pt(36)
        title_format.bold = True
        return slide
    
    # Process left side
    left_title = slide_data[left_side_key].get("title", "Left Side")
    left_text = slide_data[left_side_key].get("content", "")
    
    # Process right side
    right_title = slide_data[right_side_key].get("title", "Right Side")
    right_text = slide_data[right_side_key].get("content", "")
    
    # Split content by newlines for overflow detection
    left_paragraphs = left_text.split('\n') if left_text else []
    right_paragraphs = right_text.split('\n') if right_text else []
    
    # Determine max bullets per column
    max_per_column = MAX_BULLET_POINTS_PER_SLIDE // 2 - 1  # Subtract 1 for title
    
    # Check for overflow
    left_first, left_overflow = check_content_overflow(left_paragraphs, max_per_column)
    right_first, right_overflow = check_content_overflow(right_paragraphs, max_per_column)
    
    # Add content with titles to left side
    left_content.text_frame.clear()
    title_para = left_content.text_frame.paragraphs[0]
    title_para.text = left_title
    title_para.font.bold = True
    
    # Add content with bullet points respecting newlines
    for paragraph in left_first:
        para = left_content.text_frame.add_paragraph()
        para.text = format_text(paragraph, 30)
        para.level = 0
        para.font.size = Pt(20)
    
    # Same for right side
    right_content.text_frame.clear()
    title_para = right_content.text_frame.paragraphs[0]
    title_para.text = right_title
    title_para.font.bold = True
    
    # Add content with bullet points respecting newlines
    for paragraph in right_first:
        para = right_content.text_frame.add_paragraph()
        para.text = format_text(paragraph, 30)
        para.level = 0
        para.font.size = Pt(20)
    
    # Format title
    title_format = title.text_frame.paragraphs[0].font
    title_format.size = Pt(36)
    title_format.bold = True
    
    # Handle overflow if present
    if left_overflow or right_overflow:
        # Create continuation slide
        continuation_data = slide_data.copy()
        
        # Update content for continuation
        continuation_data[left_side_key] = continuation_data[left_side_key].copy()
        continuation_data[right_side_key] = continuation_data[right_side_key].copy()
        
        if left_overflow:
            continuation_data[left_side_key]["content"] = '\n'.join(left_overflow)
        else:
            continuation_data[left_side_key]["content"] = ""
            
        if right_overflow:
            continuation_data[right_side_key]["content"] = '\n'.join(right_overflow)
        else:
            continuation_data[right_side_key]["content"] = ""
            
        add_comparison_slide(prs, continuation_data, True)
    
    return slide

def add_content_with_caption_slide(prs, slide_data, is_continuation=False):
    """Add a slide with bullet point content and optional icon/3D image."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="content_with_caption")

    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    content_data = slide_data.get("content", {})
    if isinstance(content_data, str):
        content_data = {"title": slide_data.get("title", ""), "content": content_data}

    # Set title with appropriate continuation marker
    main_title = content_data.get("title", slide_data.get("title", ""))
    title_shape.text = f"{main_title} (cont.)" if is_continuation else main_title
    
    # Apply consistent formatting to all title paragraphs
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
    
    # Process content with line breaks at 45 characters
    content_text = content_data.get("content", "")
    formatted_content = []
    for line in content_text.split('\n'):
        words = line.split()
        current_line = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + (1 if current_length > 0 else 0) <= 45:
                current_line.append(word)
                current_length += len(word) + (1 if current_length > 0 else 0)
            else:
                formatted_content.append(" ".join(current_line))
                current_line = [word]
                current_length = len(word)
        
        if current_line:
            formatted_content.append(" ".join(current_line))
    
    # Check for overflow using MAX_BULLET_POINTS_PER_SLIDE
    first_part, overflow = check_content_overflow(formatted_content)
    
    # Add bullet points using the referenced function
    set_bullet_points(content_shape.text_frame, '\n'.join(first_part), max_line_length=45, font_size=20)

    if not is_continuation:
        # 3D icon or chart handling
        chart_data = slide_data.get("chart/smart3D_icon")
        if chart_data:
            left = Inches(5.8)
            top = Inches(2)
            width = Inches(3)
            height = Inches(2.5)
            
            # Add placeholder or actual chart/3D based on implementation
            icon_box = slide.shapes.add_textbox(left, top, width, height)
            icon_frame = icon_box.text_frame
            icon_frame.text = f"[3D Icon/Chart]"
            para = icon_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(14)

    # Recursive call for overflow
    if overflow:
        continuation_data = slide_data.copy()
        if isinstance(continuation_data.get("content", {}), dict):
            continuation_data["content"] = continuation_data["content"].copy()
            continuation_data["content"]["content"] = '\n'.join(overflow)
        else:
            continuation_data["content"] = {'title': main_title, 'content': '\n'.join(overflow)}
        add_content_with_caption_slide(prs, continuation_data, is_continuation=True)

    return slide

def add_image_with_caption_slide(prs, slide_data):
    """Add a slide with an image and caption."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_layout_styled_background(prs, slide, layout_type="image_with_caption")

    # Add title with max 60 characters, line breaks at 30 chars
    title_text = slide_data.get("title", "")
    
    # Format title with line breaks
    formatted_title_lines = []
    words = title_text.split()
    current_line = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + (1 if current_length > 0 else 0) <= 30:
            current_line.append(word)
            current_length += len(word) + (1 if current_length > 0 else 0)
        else:
            formatted_title_lines.append(" ".join(current_line))
            current_line = [word]
            current_length = len(word)
    
    if current_line:
        formatted_title_lines.append(" ".join(current_line))
    
    formatted_title = "\n".join(formatted_title_lines)
    
    # Add title textbox
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = formatted_title
    
    # Apply consistent formatting to all title paragraphs
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    # Add image
    image_path = slide_data.get("image_path", "")
    left = Inches(2)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(4)
    
    try:
        slide.shapes.add_picture(image_path, left, top, width, height)
    except Exception:
        # Add image path as placeholder text if invalid
        placeholder_box = slide.shapes.add_textbox(left, top, width, height)
        placeholder_frame = placeholder_box.text_frame
        placeholder_frame.text = f"[Image: {image_path}]"
        paragraph = placeholder_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.italic = True
        paragraph.font.size = Pt(14)

    # Add caption with max 250 characters
    caption_text = slide_data.get("content", "")
    if len(caption_text) > 250:
        caption_text = caption_text[:247] + "..."

    # Format caption with line breaks at 90 characters
    formatted_caption_lines = []
    words = caption_text.split()
    current_line = []
    current_length = 0

    for word in words:
        if current_length + len(word) + (1 if current_length > 0 else 0) <= 90:
            current_line.append(word)
            current_length += len(word) + (1 if current_length > 0 else 0)
        else:
            formatted_caption_lines.append(" ".join(current_line))
            current_line = [word]
            current_length = len(word)

    if current_line:
        formatted_caption_lines.append(" ".join(current_line))

    formatted_caption = "\n".join(formatted_caption_lines)

    caption_top = top + height + Inches(0.3)
    caption_box = slide.shapes.add_textbox(Inches(1), caption_top, Inches(8), Inches(1))
    caption_frame = caption_box.text_frame
    caption_frame.text = formatted_caption

    # Format all caption paragraphs
    for paragraph in caption_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.italic = True
        paragraph.font.size = Pt(16)

    return slide

# def add_table_slide(prs, slide_data):
#     """Add a slide with a table using the default title and dynamic sizing."""
#     slide_layout = prs.slide_layouts[5]  # Using a blank layout
#     slide = prs.slides.add_slide(slide_layout)

#     # Background styling (optional)
#     add_layout_styled_background(prs, slide, layout_type="title_with_table")

#     # Use the title placeholder if available
#     if slide.shapes.title:
#         title_shape = slide.shapes.title
#         title_shape.text = slide_data.get("title", "Table Slide")
#         for paragraph in title_shape.text_frame.paragraphs:
#             paragraph.font.size = Pt(32)
#             paragraph.font.bold = True
#             paragraph.alignment = PP_ALIGN.LEFT

#     # Get table data
#     table_data = slide_data.get("table", {})
#     headers = table_data.get("headers", [])
#     rows = table_data.get("rows", [])

#     if not headers or not rows:
#         return slide

#     n_rows = len(rows) + 1
#     n_cols = len(headers)

#     # Set table position
#     left = Inches(0.5)
#     top = Inches(1.5)
#     max_width = Inches(9)
#     max_height = Inches(5)

#     # Estimate column widths based on max string length
#     max_lengths = [len(str(header)) for header in headers]
#     for row in rows:
#         for i, cell in enumerate(row):
#             max_lengths[i] = max(max_lengths[i], len(str(cell)))

#     total_length = sum(max_lengths)
#     col_widths = [max_width * (l / total_length) for l in max_lengths]

#     # Add table with approximate dimensions
#     table = slide.shapes.add_table(n_rows, n_cols, left, top, max_width, max_height).table

#     # Apply calculated widths
#     for i, width in enumerate(col_widths):
#         table.columns[i].width = int(width)

#     # Add headers
#     for col_idx, header in enumerate(headers):
#         cell = table.cell(0, col_idx)
#         cell.text = str(header)
#         for para in cell.text_frame.paragraphs:
#             para.font.bold = True
#             para.font.size = Pt(14)
#             para.alignment = PP_ALIGN.CENTER
#         cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

#     # Fill rows
#     for row_idx, row_data in enumerate(rows):
#         for col_idx, cell_data in enumerate(row_data):
#             cell = table.cell(row_idx + 1, col_idx)
#             cell.text = str(cell_data)
#             for para in cell.text_frame.paragraphs:
#                 para.font.size = Pt(13)
#                 para.alignment = PP_ALIGN.CENTER
#             cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

#     # Dynamic row heights based on number of lines
#     for i in range(n_rows):
#         row_height = Inches(0.4)
#         for j in range(n_cols):
#             lines = table.cell(i, j).text.count('\n') + 1
#             row_height = max(row_height, Inches(0.2 + 0.2 * lines))
#         table.rows[i].height = int(row_height)

#     # Optional caption/content below the table
#     content_text = slide_data.get("content", "")
#     if content_text:
#         content_top = top + max_height + Inches(0.3)
#         content_box = slide.shapes.add_textbox(left, content_top, max_width, Inches(2))
#         content_frame = content_box.text_frame
#         content_frame.word_wrap = True
#         content_frame.clear()  # Clear default paragraph

#         for line in content_text.strip().split('\n'):
#             p = content_frame.add_paragraph()
#             p.text = line
#             p.font.size = Pt(12)
#             p.alignment = PP_ALIGN.LEFT

#     return slide

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

# def add_table_slide(prs, slide_data):
#     """
#     Add one or more slides with a table, automatically handling pagination.
#     Parameters:
#         prs: PowerPoint presentation object
#         slide_data: Dictionary containing slide title and table data
#     Returns:
#         List of created slides
#     """
#     table_data = slide_data.get("table", {})
#     headers = table_data.get("headers", [])
#     rows = table_data.get("rows", [])

#     if not headers or not rows:
#         return [create_empty_table_slide(prs, slide_data)]

#     # Calculate max visible rows per slide to prevent overflow
#     # More conservative calculation to ensure no overflow
#     max_rows = calculate_max_rows_for_slide(headers, rows)
#     slides = []
#     total_pages = (len(rows) + max_rows - 1) // max_rows if max_rows > 0 else 1

#     for i in range(0, len(rows), max_rows):
#         chunk_rows = rows[i:i + max_rows]
#         page_num = (i // max_rows) + 1
        
#         page_title = slide_data.get('title', 'Table')
#         if total_pages > 1:
#             page_title = f"{page_title} ({page_num}/{total_pages})"
            
#         page_data = {
#             **slide_data,
#             "title": page_title,
#             "table": {"headers": headers, "rows": chunk_rows},
#             "page_number": page_num,
#             "total_pages": total_pages
#         }

#         slides.append(create_table_slide(prs, page_data))

#     return slides

# def create_table_slide(prs, slide_data):
#     """
#     Create a slide with a formatted table using the proper title placeholder
#     """
#     # Use slide layout 5 (Title and Content) instead of blank layout
#     slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
#     apply_slide_background(slide)

#     # Use the built-in title placeholder
#     if slide.shapes.title:
#         title_shape = slide.shapes.title
#         title_shape.text = slide_data.get("title", "Table")
        
#         # Style the title
#         for paragraph in title_shape.text_frame.paragraphs:
#             paragraph.font.size = Pt(32)
#             paragraph.font.bold = True
#             paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
#     # Clear any existing content placeholder to make room for our custom table
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape != slide.shapes.title:
#             if hasattr(shape, 'placeholder_format'):
#                 if shape.placeholder_format.type == 7:  # Content placeholder
#                     shape.element.getparent().remove(shape.element)

#     # Calculate available space for table
#     # Title takes up about 1 inch vertical space
#     title_height = Inches(1.0)
#     footer_height = Inches(0.4)  # Space for footer
    
#     # Available height for table (7.5 is slide height)
#     available_height = Inches(7.5) - title_height - footer_height
    
#     # Position table right below the title
#     table_top = title_height
    
#     # Create table with calculated dimensions
#     table_data = slide_data.get("table", {})
#     headers = table_data.get("headers", [])
#     rows = table_data.get("rows", [])

#     n_rows, n_cols = len(rows) + 1, len(headers)
#     table_shape = slide.shapes.add_table(
#         n_rows, n_cols, Inches(0.5), table_top, Inches(9), available_height
#     )
#     table = table_shape.table

#     # Style header row
#     for i, header in enumerate(headers):
#         cell = table.cell(0, i)
#         # Use shorter text for overflow prevention
#         cell.text = format_cell_text_strict(str(header), 25)
#         style_header_cell(cell)

#     # Style data rows
#     for i, row in enumerate(rows):
#         for j, val in enumerate(row):
#             if j < n_cols:  # Ensure we don't go out of bounds
#                 cell = table.cell(i + 1, j)
#                 # Use shorter text for overflow prevention
#                 cell.text = format_cell_text_strict(str(val), 25)
#                 style_data_cell(cell)

#     # Optimize column widths and row heights with strict height limits
#     optimize_table_layout_strict(table, headers, rows, available_height / (n_rows))
    
#     # Add footer with page numbers if multiple pages
#     if slide_data.get("total_pages", 1) > 1:
#         add_page_footer(slide, slide_data.get("page_number", 1), slide_data.get("total_pages", 1))

#     return slide

# def create_empty_table_slide(prs, slide_data):
#     """Create a slide indicating no data is available"""
#     # Use Title and Content layout
#     slide = prs.slides.add_slide(prs.slide_layouts[1])
#     apply_slide_background(slide)
    
#     # Use built-in title placeholder
#     if slide.shapes.title:
#         title_shape = slide.shapes.title
#         title_shape.text = slide_data.get("title", "Table")
        
#         # Style the title
#         for paragraph in title_shape.text_frame.paragraphs:
#             paragraph.font.size = Pt(32)
#             paragraph.font.bold = True
#             paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
#     # Add no data message
#     msg_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
#     tf = msg_shape.text_frame
#     p = tf.add_paragraph()
#     p.text = "No table data available"
#     p.font.size = Pt(18)
#     p.font.italic = True
#     p.alignment = PP_ALIGN.CENTER
    
#     return slide

# def calculate_max_rows_for_slide(headers, rows):
#     """
#     Calculate maximum rows that can fit on a slide without overflow
#     Using more conservative estimates
#     """
#     # Fixed heights in presentation
#     title_height = 1.0  # inches
#     footer_height = 0.4  # inches
#     available_vertical_space = 7.5 - title_height - footer_height  # 7.5 is standard slide height
    
#     # Header takes fixed space
#     header_row_height = 0.4  # inches
#     available_for_data = available_vertical_space - header_row_height
    
#     # Analyze sample rows to determine average row height
#     sample_size = min(5, len(rows))
#     if sample_size == 0:
#         return 10  # Default if no rows
        
#     # Estimate content density more accurately
#     max_chars_per_line = 25  # More strict limit for text width
#     max_lines_per_cell = []
    
#     for i in range(sample_size):
#         row_lines = []
#         for cell in rows[i]:
#             cell_text = str(cell)
#             # Calculate how many lines this would take with strict wrapping
#             char_count = len(cell_text)
#             if char_count == 0:
#                 lines_needed = 1
#             else:
#                 lines_needed = (char_count // max_chars_per_line) + (1 if char_count % max_chars_per_line > 0 else 0)
#                 lines_needed = max(1, lines_needed)
#             row_lines.append(lines_needed)
#         max_lines_per_cell.append(max(row_lines))
    
#     avg_max_lines = sum(max_lines_per_cell) / len(max_lines_per_cell)
    
#     # Calculate space needed per row with padding
#     # Each line of text takes about 0.2 inches, plus cell padding
#     row_height = 0.2 + (0.15 * avg_max_lines)
    
#     # Add 10% safety margin
#     row_height *= 1.1
    
#     # Calculate max rows that fit, ensuring minimum of at least 3 rows
#     max_rows = max(3, int(available_for_data / row_height))
    
#     # Cap at reasonable limits - never more than 12 rows per slide to avoid overflows
#     return min(12, max_rows)

# def format_cell_text_strict(text, max_chars=25):
#     """Format text for table cells with strict line breaks to prevent overflow"""
#     if not text:
#         return ""
        
#     # For very short text, return as is
#     if len(text) <= max_chars:
#         return text
    
#     # For longer text, apply strict wrapping
#     words = text.split()
#     lines = []
#     current_line = []
#     current_length = 0
    
#     for word in words:
#         # Check if adding this word would exceed the line length
#         if current_length + len(word) + (1 if current_length > 0 else 0) > max_chars:
#             # If current line has content, add it to lines
#             if current_line:
#                 lines.append(" ".join(current_line))
#                 current_line = []
#                 current_length = 0
            
#             # Handle long words by truncating if necessary
#             if len(word) > max_chars:
#                 # Break the word into chunks
#                 while word:
#                     chunk = word[:max_chars-1] + "-" if len(word) > max_chars-1 else word
#                     lines.append(chunk)
#                     word = word[max_chars-1:] if len(word) > max_chars-1 else ""
#             else:
#                 current_line = [word]
#                 current_length = len(word)
#         else:
#             current_line.append(word)
#             current_length += len(word) + (1 if current_length > 0 else 0)
    
#     # Add any remaining content
#     if current_line:
#         lines.append(" ".join(current_line))
    
#     # Limit total lines to prevent excessive vertical space
#     if len(lines) > 5:
#         lines = lines[:4]
#         lines.append("...")
    
#     return "\n".join(lines)

# def style_header_cell(cell):
#     """Apply specific styling to header cells"""
#     tf = cell.text_frame
#     tf.vertical_anchor = MSO_ANCHOR.MIDDLE
#     tf.word_wrap = True
    
#     # Set cell background
#     cell.fill.solid()
#     cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
    
#     # Apply text styles
#     for paragraph in tf.paragraphs:
#         paragraph.alignment = PP_ALIGN.CENTER
#         paragraph.font.bold = True
#         paragraph.font.size = Pt(12)
#         paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
#     # Set margins for better readability
#     cell.margin_left = Inches(0.05)
#     cell.margin_right = Inches(0.05)
#     cell.margin_top = Inches(0.03)
#     cell.margin_bottom = Inches(0.03)

# def style_data_cell(cell):
#     """Apply styling to data cells"""
#     tf = cell.text_frame
#     tf.vertical_anchor = MSO_ANCHOR.MIDDLE
#     tf.word_wrap = True
    
#     # Apply alternating row colors if needed
#     # row_idx = cell.row_idx
#     # if row_idx % 2 == 1:  # Odd rows
#     #     cell.fill.solid()
#     #     cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
#     # Apply text styles
#     for paragraph in tf.paragraphs:
#         paragraph.alignment = PP_ALIGN.LEFT
#         paragraph.font.size = Pt(11)
#         paragraph.font.color.rgb = RGBColor(50, 50, 50)
    
#     # Set margins for better readability
#     cell.margin_left = Inches(0.05)
#     cell.margin_right = Inches(0.05)
#     cell.margin_top = Inches(0.03)
#     cell.margin_bottom = Inches(0.03)

# def optimize_table_layout_strict(table, headers, rows, max_row_height):
#     """
#     Optimize column widths and row heights with strict height limits
#     """
#     # Calculate optimal column widths based on content
#     col_widths = calculate_column_widths(headers, rows)
#     total_width = Inches(9)  # Total available width
    
#     # Apply calculated widths while ensuring minimum sizes
#     for i, width_percent in enumerate(col_widths):
#         if i < len(table.columns):
#             min_width = Inches(0.5)  # Minimum column width
#             table.columns[i].width = max(min_width, int(total_width * width_percent))
    
#     # Set header row height
#     if table.rows:
#         table.rows[0].height = Inches(0.4)
    
#     # Optimize data row heights with strict limits
#     for i in range(1, len(table.rows)):
#         # Count lines in the row
#         max_lines = 1
#         for j in range(len(table.columns)):
#             if i < len(table.rows) and j < len(table.columns):
#                 cell_text = table.cell(i, j).text
#                 lines_count = cell_text.count('\n') + 1
#                 max_lines = max(max_lines, lines_count)
        
#         # Calculate row height based on content
#         content_based_height = Inches(0.2 + (0.15 * max_lines))
        
#         # Apply upper limit to prevent overflow
#         final_height = min(content_based_height, Inches(max_row_height))
        
#         # Set row height
#         if i < len(table.rows):
#             table.rows[i].height = max(Inches(0.25), final_height)

# def calculate_column_widths(headers, rows):
#     """Calculate balanced column widths with content awareness"""
#     # Get max character length for each column
#     col_lengths = []
    
#     # Initialize with header lengths
#     for header in headers:
#         col_lengths.append(len(str(header)))
    
#     # Check data rows (sample for efficiency)
#     sample_size = min(20, len(rows))
#     for i in range(sample_size):
#         for j, val in enumerate(rows[i]):
#             if j < len(col_lengths):
#                 col_lengths[j] = max(col_lengths[j], len(str(val)))
    
#     # Calculate percentages with minimum widths
#     total_length = sum(col_lengths)
#     min_percent = 0.1  # Minimum 10% of total width
    
#     if total_length == 0:
#         # Equal distribution if no content
#         return [1.0 / len(headers)] * len(headers)
    
#     # Calculate initial percentages
#     percentages = [length / total_length for length in col_lengths]
    
#     # Balance percentages - no column should be too dominant
#     max_percent = 0.3  # Maximum 30% of total width
#     for i in range(len(percentages)):
#         if percentages[i] > max_percent:
#             excess = percentages[i] - max_percent
#             percentages[i] = max_percent
#             # Distribute excess to other columns proportionally
#             total_others = sum(percentages) - max_percent
#             if total_others > 0:
#                 for j in range(len(percentages)):
#                     if j != i:
#                         percentages[j] += excess * (percentages[j] / total_others)
    
#     # Apply minimum width constraints
#     for i in range(len(percentages)):
#         if percentages[i] < min_percent:
#             percentages[i] = min_percent
    
#     # Normalize to ensure sum is 1.0
#     total_percent = sum(percentages)
#     if total_percent > 0:
#         percentages = [p / total_percent for p in percentages]
    
#     return percentages

# def add_page_footer(slide, page_num, total_pages):
#     """Add page number footer to slide"""
#     footer = slide.shapes.add_textbox(
#         Inches(0.5), Inches(7.0), Inches(9), Inches(0.3)
#     )
#     tf = footer.text_frame
#     p = tf.add_paragraph()
#     p.text = f"Page {page_num} of {total_pages}"
#     p.font.size = Pt(10)
#     p.font.color.rgb = RGBColor(100, 100, 100)
#     p.alignment = PP_ALIGN.RIGHT

# def apply_slide_background(slide):
#     """Apply a subtle gradient background to the slide"""
#     fill = slide.background.fill
#     fill.gradient()
#     fill.gradient_angle = 90
#     fill.gradient_stops[0].position = 0
#     fill.gradient_stops[0].color.rgb = RGBColor(250, 250, 252)
#     fill.gradient_stops[1].position = 1
#     fill.gradient_stops[1].color.rgb = RGBColor(240, 245, 250)

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

def add_table_slide(prs, slide_data):
    """
    Add one or more slides with a table, automatically handling pagination.
    Parameters:
        prs: PowerPoint presentation object
        slide_data: Dictionary containing slide title and table data
    Returns:
        List of created slides
    """
    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers or not rows:
        return [create_empty_table_slide(prs, slide_data)]

    # Use adaptive calculation to maximize slide space usage
    max_rows = calculate_max_rows_adaptive(headers, rows)
    slides = []
    total_pages = (len(rows) + max_rows - 1) // max_rows if max_rows > 0 else 1

    for i in range(0, len(rows), max_rows):
        chunk_rows = rows[i:i + max_rows]
        page_num = (i // max_rows) + 1
        
        page_title = slide_data.get('title', 'Table')
        if total_pages > 1:
            page_title = f"{page_title} ({page_num}/{total_pages})"
            
        page_data = {
            **slide_data,
            "title": page_title,
            "table": {"headers": headers, "rows": chunk_rows},
            "page_number": page_num,
            "total_pages": total_pages
        }

        slides.append(create_table_slide(prs, page_data))

    return slides

def create_table_slide(prs, slide_data):
    """
    Create a slide with a formatted table using the proper title placeholder
    """
    # Use slide layout 1 (Title and Content)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_background(slide)

    # Use the built-in title placeholder
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Table")
        
        # Style the title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    # Clear any existing content placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            if hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 7:  # Content placeholder
                    shape.element.getparent().remove(shape.element)

    # Calculate available space more accurately
    title_height = Inches(0.95)  # Slightly reduced from previous 1.0
    footer_height = Inches(0.3)   # Reduced from previous 0.4
    
    # Standard slide height is 7.5 inches
    available_height = Inches(7.5) - title_height - footer_height
    
    # Position table right below the title
    # table_top = title_height
    if slide.shapes.title:
        actual_title_height = slide.shapes.title.height
        table_top = slide.shapes.title.top + actual_title_height + Inches(0.15)  # Add 0.15 inch buffer
    else:
        table_top = title_height
    
    # Create table with calculated dimensions
    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    n_rows, n_cols = len(rows) + 1, len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), table_top, Inches(9), available_height
    )
    table = table_shape.table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = format_cell_text_adaptive(str(header), 25)
        style_header_cell(cell)

    # Style data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j < n_cols:  # Ensure we don't go out of bounds
                cell = table.cell(i + 1, j)
                cell.text = format_cell_text_adaptive(str(val), 25)
                style_data_cell(cell)

    # Optimize layout to maximize space usage
    optimize_table_layout_adaptive(table, headers, rows, n_rows)
    
    # Add footer with page numbers if multiple pages
    if slide_data.get("total_pages", 1) > 1:
        add_page_footer(slide, slide_data.get("page_number", 1), slide_data.get("total_pages", 1))

    return slide

def create_empty_table_slide(prs, slide_data):
    """Create a slide indicating no data is available"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_slide_background(slide)
    
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "Table")
        
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    msg_shape = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf = msg_shape.text_frame
    p = tf.add_paragraph()
    p.text = "No table data available"
    p.font.size = Pt(18)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def calculate_max_rows_adaptive(headers, rows):
    """
    Calculate maximum rows based on actual content analysis
    with minimal safety margins to maximize slide usage
    """
    # Standard measurements
    slide_height = 7.5  # inches
    title_height = 0.95  # inches
    footer_height = 0.3  # inches
    header_row_height = 0.4  # inches
    
    # Calculate usable height for data rows
    usable_height = slide_height - title_height - footer_height - header_row_height
    
    # Analyze actual content to make better decisions
    if not rows:
        return 15  # Default if no rows
    
    # Sample rows for analysis (more samples for better accuracy)
    sample_size = min(10, len(rows))
    row_heights = []
    
    for i in range(sample_size):
        row = rows[i] if i < len(rows) else rows[-1]
        # Analyze each cell in the row
        max_lines_in_row = 1
        for cell_content in row:
            text = str(cell_content)
            # Estimate lines based on reasonable word wrapping
            if len(text) <= 25:  # Short text
                lines = 1
            else:
                # Approximate line count based on text length and avg words per line
                words = text.split()
                avg_word_length = sum(len(word) for word in words) / len(words) if words else 5
                chars_per_line = 25  # Target chars per line
                words_per_line = chars_per_line / (avg_word_length + 1)  # +1 for space
                lines = max(1, len(words) / words_per_line if words_per_line > 0 else 1)
                
                # Cap at 6 lines per cell for display
                lines = min(6, lines)
            
            max_lines_in_row = max(max_lines_in_row, lines)
        
        # Calculate estimated row height (base + line height × line count)
        estimated_height = 0.2 + (0.13 * max_lines_in_row)  # Slightly reduced line height
        row_heights.append(estimated_height)
    
    # Calculate average row height - bias toward taller rows for safety
    if row_heights:
        # Use 75th percentile instead of average for better accuracy
        row_heights.sort()
        percentile_idx = int(len(row_heights) * 0.75)
        typical_row_height = row_heights[percentile_idx]
    else:
        typical_row_height = 0.3  # Default if calculation fails
    
    # Add minimal safety margin (5%)
    typical_row_height *= 1.05
    
    # Calculate how many rows fit
    max_rows = int(usable_height / typical_row_height)
    
    # Set reasonable bounds - at least 5 rows, at most 20
    return max(5, min(20, max_rows))

def format_cell_text_adaptive(text, max_chars=25):
    """Format text for table cells with adaptive line breaks"""
    if not text:
        return ""
        
    # For short text, return as is
    if len(text) <= max_chars:
        return text
    
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        # Check if adding this word would exceed the line length
        if current_length + len(word) + (1 if current_length > 0 else 0) > max_chars:
            if current_line:
                lines.append(" ".join(current_line))
                current_line = []
                current_length = 0
            
            # Handle long words
            if len(word) > max_chars:
                # Break the word into chunks
                while word:
                    chunk = word[:max_chars-1] + "-" if len(word) > max_chars-1 else word
                    lines.append(chunk)
                    word = word[max_chars-1:] if len(word) > max_chars-1 else ""
            else:
                current_line = [word]
                current_length = len(word)
        else:
            current_line.append(word)
            current_length += len(word) + (1 if current_length > 0 else 0)
    
    # Add any remaining content
    if current_line:
        lines.append(" ".join(current_line))
    
    # Limit total lines to prevent excessive vertical space but allow more lines
    # to maximize space utilization (increased from 5 to 6)
    if len(lines) > 6:
        lines = lines[:5]
        lines.append("...")
    
    return "\n".join(lines)

def style_header_cell(cell):
    """Apply specific styling to header cells"""
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    
    # Set cell background
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
    
    # Apply text styles
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    # Set margins - reduced slightly for more compact layout
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)

def style_data_cell(cell):
    """Apply styling to data cells with space efficiency"""
    tf = cell.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    
    # Apply text styles
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = RGBColor(50, 50, 50)
    
    # Set compact margins to maximize content space
    cell.margin_left = Inches(0.04)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)

def optimize_table_layout_adaptive(table, headers, rows, total_rows):
    """
    Optimize column widths and row heights to maximize content display
    """
    # Calculate optimal column widths based on content
    col_widths = calculate_column_widths_adaptive(headers, rows)
    total_width = Inches(9)  # Total available width
    
    # Apply calculated widths
    for i, width_percent in enumerate(col_widths):
        if i < len(table.columns):
            min_width = Inches(0.4)  # Reduced minimum width
            table.columns[i].width = max(min_width, int(total_width * width_percent))
    
    # Set header row height - compact but readable
    if table.rows:
        table.rows[0].height = Inches(0.35)  # Slightly reduced
    
    # Optimize data row heights to fit content efficiently
    for i in range(1, len(table.rows)):
        # Calculate lines in this row's cells
        max_lines = 1
        for j in range(len(table.columns)):
            if i < len(table.rows) and j < len(table.columns):
                cell_text = table.cell(i, j).text
                lines_count = cell_text.count('\n') + 1
                max_lines = max(max_lines, lines_count)
        
        # More efficient row height calculation
        # Base height + height per line, with reduced values
        row_height = Inches(0.18 + (0.13 * max_lines))
        
        # Ensure minimum readability
        min_height = Inches(0.24)
        
        if i < len(table.rows):
            table.rows[i].height = max(min_height, row_height)

def calculate_column_widths_adaptive(headers, rows):
    """Calculate balanced column widths based on actual content needs"""
    # Initialize column content analysis
    col_data = []
    for _ in range(len(headers)):
        col_data.append({
            'max_length': 0,       # Max character length
            'avg_length': 0,       # Average character length
            'long_word_count': 0,  # Number of long words (>10 chars)
            'sample_count': 0      # Number of samples analyzed
        })
    
    # Analyze headers
    for i, header in enumerate(headers):
        text = str(header)
        col_data[i]['max_length'] = max(col_data[i]['max_length'], len(text))
        col_data[i]['avg_length'] += len(text)
        col_data[i]['sample_count'] += 1
        
        # Count long words
        words = text.split()
        for word in words:
            if len(word) > 10:
                col_data[i]['long_word_count'] += 1
    
    # Analyze sample rows
    sample_size = min(20, len(rows))
    for i in range(sample_size):
        for j, val in enumerate(rows[i]):
            if j < len(col_data):
                text = str(val)
                col_data[j]['max_length'] = max(col_data[j]['max_length'], len(text))
                col_data[j]['avg_length'] += len(text)
                col_data[j]['sample_count'] += 1
                
                # Count long words
                words = text.split()
                for word in words:
                    if len(word) > 10:
                        col_data[j]['long_word_count'] += 1
    
    # Calculate average lengths
    for data in col_data:
        if data['sample_count'] > 0:
            data['avg_length'] /= data['sample_count']
    
    # Calculate importance score for each column
    importance_scores = []
    for data in col_data:
        # Formula considers max length, average length, and presence of long words
        score = (0.4 * data['max_length'] + 
                0.4 * data['avg_length'] + 
                0.2 * data['long_word_count'])
        importance_scores.append(max(1, score))  # Ensure minimum score of 1
    
    # Calculate percentages based on importance scores
    total_score = sum(importance_scores)
    if total_score == 0:
        # Equal distribution if no content
        return [1.0 / len(headers)] * len(headers)
    
    # Initial percentages based on importance
    percentages = [score / total_score for score in importance_scores]
    
    # Apply reasonable constraints (min 7%, max 30%)
    min_percent = 0.07
    max_percent = 0.30
    
    for i in range(len(percentages)):
        if percentages[i] < min_percent:
            percentages[i] = min_percent
        elif percentages[i] > max_percent:
            percentages[i] = max_percent
    
    # Normalize to ensure sum is 1.0
    total_percent = sum(percentages)
    if total_percent > 0:
        percentages = [p / total_percent for p in percentages]
    
    return percentages

def add_page_footer(slide, page_num, total_pages):
    """Add compact page number footer to slide"""
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.1), Inches(9), Inches(0.2)
    )
    tf = footer.text_frame
    p = tf.add_paragraph()
    p.text = f"Page {page_num} of {total_pages}"
    p.font.size = Pt(9)  # Smaller font
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.RIGHT

def apply_slide_background(slide):
    """Apply a subtle gradient background to the slide"""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 90
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = RGBColor(250, 250, 252)
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(240, 245, 250)


def format_text(text, max_line_length):
    """Format text to wrap at specified character length, keeping it as a single paragraph."""
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        if current_length + len(word) + len(current_line) <= max_line_length:
            current_line.append(word)
            current_length += len(word)
        else:
            lines.append(' '.join(current_line))
            current_line = [word]
            current_length = len(word)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    # Join the lines with spaces rather than newlines to keep in same paragraph
    return ' '.join(lines)