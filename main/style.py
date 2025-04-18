from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

def add_layout_styled_background(prs, slide, layout_type):
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_background_fill(color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*color)
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def add_center_card(color, opacity=200):
        card_width = int(slide_width)
        card_height = int(slide_height * 0.75)
        card_left = int((slide_width - card_width) / 2)
        card_top = int((slide_height - card_height) / 1)

        card = slide.shapes.add_shape(
            # MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height
            MSO_SHAPE.RECTANGLE, card_left, card_top, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*color)
        card.fill.fore_color.alpha = opacity
        card.line.fill.background()
        slide.shapes._spTree.remove(card._element)
        slide.shapes._spTree.insert(3, card._element)

    def add_bottom_stripe(color, height_ratio=0.1):
        stripe_height = int(slide_height * height_ratio)
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, slide_height - stripe_height, slide_width, stripe_height
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(*color)
        stripe.line.fill.background()
        slide.shapes._spTree.remove(stripe._element)
        slide.shapes._spTree.insert(3, stripe._element)

    def add_side_cards(left_color, right_color, opacity=180):
        card_width = int(slide_width * 0.45)
        card_height = int(slide_height * 0.75)
        padding = int(slide_width * 0.04)
        top = int(slide_height * 0.20)

        for i, color in enumerate([left_color, right_color]):
            left = padding + i * (card_width + padding)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*color)
            card.fill.fore_color.alpha = opacity
            card.line.fill.background()
            slide.shapes._spTree.remove(card._element)
            slide.shapes._spTree.insert(3, card._element)

    def add_table_slide_style(prs, slide, layout_type="title_with_table"):
        """
        Apply a styled background for table slides that matches other layout styles
        
        Parameters:
            prs: PowerPoint presentation object
            slide: The slide to apply styling to
        """
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Create light background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(234, 239, 242)  # Light color
        bg.line.fill.background()
        
        # Make sure background is behind other elements
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
        
        # Add thin accent bar on top for visual interest
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            0, 
            int(slide_height * 0.05),  # Positioned below title
            slide_width, 
            int(slide_height * 0.01)   # Very thin bar
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(39, 68, 114)  # Dark accent color
        accent_bar.line.fill.background()
        slide.shapes._spTree.remove(accent_bar._element)
        slide.shapes._spTree.insert(3, accent_bar._element)
        
        # Add bottom stripe similar to title_and_content layout
        stripe_height = int(slide_height * 0.08)  # Slightly thinner than original
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            0, 
            slide_height - stripe_height, 
            slide_width, 
            stripe_height
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(39, 68, 114)  # Dark accent color
        stripe.line.fill.background()
        slide.shapes._spTree.remove(stripe._element)
        slide.shapes._spTree.insert(3, stripe._element)
    
    def add_top_accent_bar(color, height_ratio=0.01, position_ratio=0.05):
        bar_height = int(slide_height * height_ratio)
        bar_top = int(slide_height * position_ratio)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, bar_top, slide_width, bar_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.fill.background()
        slide.shapes._spTree.remove(bar._element)
        slide.shapes._spTree.insert(3, bar._element)

    # Style each layout uniquely
    if layout_type == "title_slide":
        add_background_fill((39, 68, 114))  # light_color
        add_center_card((234, 239, 242), 200)

    elif layout_type == "title_only":
        add_background_fill((234, 239, 242))

    elif layout_type == "title_and_content":
        add_background_fill((234, 239, 242))  # Slightly lighter
        add_bottom_stripe((39, 68, 114)) # dark color here

    elif layout_type == "section_header":
        add_background_fill((234, 239, 242))
        top_half = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height // 2)
        top_half.fill.solid()
        top_half.fill.fore_color.rgb = RGBColor(39, 68, 114)
        top_half.line.fill.background()
        slide.shapes._spTree.remove(top_half._element)
        slide.shapes._spTree.insert(3, top_half._element)

    elif layout_type == "two_content":
        add_background_fill((39, 68, 114))
        add_side_cards((234, 239, 242), (234, 239, 242), opacity=180)

    elif layout_type == "comparison":
        add_background_fill((39, 68, 114))
        add_side_cards((234, 239, 242), (234, 239, 242), opacity=180)

    elif layout_type in ["content_with_caption", "image_with_caption"]:
        add_background_fill((234, 239, 242))

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            int(slide_height * 0.15),
            int(slide_width * 0.02),
            int(slide_height * 0.7)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(39, 68, 114)
        bar.line.fill.background()
        slide.shapes._spTree.remove(bar._element)
        slide.shapes._spTree.insert(3, bar._element)

        # 🔧 Move and resize text box
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text.strip() != "":
                shape.left = int(slide_width * 0.06)
                shape.top = int(slide_height * 0.2)
                shape.width = int(slide_width * 0.5)
                shape.height = int(slide_height * 0.6)

        # 🔧 Move and resize image box
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape.left = int(slide_width * 0.6)
                shape.top = int(slide_height * 0.25)
                shape.width = int(slide_width * 0.35)
                shape.height = int(slide_height * 0.5)

    elif layout_type == "table_slide":
        # Light background like the title_and_content layout
        add_background_fill((234, 239, 242))
        
        # Add thin accent bar below title area
        add_top_accent_bar((39, 68, 114))
        
        # Add bottom stripe similar to title_and_content layout but slightly thinner
        add_bottom_stripe((39, 68, 114), height_ratio=0.08)

    else:
        add_background_fill((39, 68, 114))