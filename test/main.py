import os
import json
import argparse
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

def log(msg):
    print(f"[LOG] {msg}")

def add_text_to_shape(shape, text, font_size=None, font_name=None, bold=None, italic=None, color=None):
    """Add text to a shape with formatting options"""
    if not hasattr(shape, "text_frame"):
        log(f"⚠️ Shape doesn't have a text frame")
        return
    
    # Clear existing text
    shape.text_frame.clear()
    
    # Handle multi-line text
    paragraphs = text.split('\n')
    
    # Add first paragraph
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = paragraphs[0]
    
    # Apply formatting to first run
    if font_size:
        run.font.size = Pt(font_size)
    if font_name:
        run.font.name = font_name
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color:
        try:
            run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
        except Exception as e:
            log(f"Color parsing error: {e}")
    
    # Add additional paragraphs if needed
    for para_text in paragraphs[1:]:
        p = shape.text_frame.add_paragraph()
        run = p.add_run()
        run.text = para_text
        
        # Apply same formatting to additional paragraphs
        if font_size:
            run.font.size = Pt(font_size)
        if font_name:
            run.font.name = font_name
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if color:
            try:
                run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
            except Exception as e:
                log(f"Color parsing error: {e}")

def add_image_to_placeholder(slide, placeholder, image_path):
    """Add image to a placeholder"""
    if not os.path.exists(image_path):
        log(f"⚠️ Image not found: {image_path}")
        return False
    
    try:
        # Calculate position and size
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        
        # Add picture
        slide.shapes.add_picture(image_path, left, top, width, height)
        return True
    except Exception as e:
        log(f"⚠️ Error adding image: {e}")
        return False

def find_placeholder_by_type(slide, placeholder_type):
    """Find a placeholder in a slide by its type"""
    for shape in slide.shapes:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            if shape.placeholder_format.type == placeholder_type:
                return shape
    return None

def find_placeholder_by_index(slide, idx):
    """Find a placeholder in a slide by its index"""
    for shape in slide.shapes:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            if shape.placeholder_format.idx == idx:
                return shape
    return None

def build_enhanced_presentation(schema_path, content_path, output_path, image_dir="assets"):
    """Build a PowerPoint presentation using a layout schema and content JSON"""
    if not os.path.exists(schema_path):
        log("❌ Schema file not found.")
        return False
    
    if not os.path.exists(content_path):
        log("❌ Content file not found.")
        return False
    
    # Ensure image directory exists
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)
        log(f"Created image directory: {image_dir}")
    
    # Load schema and content
    with open(schema_path, "r") as f:
        schema = json.load(f)
    
    with open(content_path, "r", encoding="utf-8") as f:
        content = json.load(f)
    
    # Create presentation
    prs = Presentation()
    
    # Process each slide
    for slide_idx, slide_content in enumerate(content):
        # Get layout type
        layout_type = slide_content.get("layout")
        if not layout_type:
            log(f"⚠️ Slide {slide_idx+1} has no layout specified. Skipping.")
            continue
        
        # Find layout in schema
        layout_info = None
        for layout in schema:
            if layout.get("layout") == layout_type:
                layout_info = layout
                break
        
        if not layout_info:
            log(f"⚠️ Layout '{layout_type}' not found in schema. Using blank slide.")
            # Use blank layout
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        else:
            layout_index = layout_info.get("layout_index", 0)
            if layout_index >= len(prs.slide_layouts):
                log(f"⚠️ Layout index {layout_index} out of range. Using default layout.")
                slide = prs.slides.add_slide(prs.slide_layouts[0])
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        log(f"✅ Created slide {slide_idx+1} with layout: {layout_type}")

        # Map common content keys to placeholder types
        content_mapping = {
            "title": PP_PLACEHOLDER.TITLE,
            "subtitle": PP_PLACEHOLDER.SUBTITLE,
            "content": PP_PLACEHOLDER.BODY,
            "body": PP_PLACEHOLDER.BODY,
            "left_content": PP_PLACEHOLDER.OBJECT,  # First content placeholder
            "right_content": PP_PLACEHOLDER.OBJECT,  # Second content placeholder
            "left_heading": PP_PLACEHOLDER.TEXT,
            "right_heading": PP_PLACEHOLDER.TEXT,
            "image": PP_PLACEHOLDER.PICTURE
        }
        
        # Process each content item
        for content_key, content_value in slide_content.items():
            if content_key == "layout":
                continue
            
            # Handle special case for notes
            if content_key == "notes":
                if hasattr(slide, "notes_slide"):
                    slide.notes_slide.notes_text_frame.text = content_value
                continue
            
            # Handle special case for images
            if content_key == "image" and isinstance(content_value, str):
                image_placeholder = find_placeholder_by_type(slide, PP_PLACEHOLDER.PICTURE)
                if image_placeholder:
                    image_path = os.path.join(image_dir, content_value)
                    add_image_to_placeholder(slide, image_placeholder, image_path)
                    log(f"Added image: {content_value}")
                continue
            
            # Try to find matching placeholder
            placeholder_type = content_mapping.get(content_key)
            
            if placeholder_type:
                # For left/right content, we need special handling
                if content_key == "left_content":
                    shape = find_placeholder_by_index(slide, 1)  # Typically first body placeholder
                elif content_key == "right_content":
                    shape = find_placeholder_by_index(slide, 2)  # Typically second body placeholder
                elif content_key == "left_heading":
                    shape = find_placeholder_by_index(slide, 3)  # Typically first text placeholder
                elif content_key == "right_heading":
                    shape = find_placeholder_by_index(slide, 4)  # Typically second text placeholder
                else:
                    shape = find_placeholder_by_type(slide, placeholder_type)
                
                if shape:
                    add_text_to_shape(shape, content_value, font_size=None)
                    log(f"Added '{content_key}' to slide {slide_idx+1}")
                else:
                    log(f"⚠️ No placeholder found for '{content_key}'")
    
    # Save presentation
    prs.save(output_path)
    log(f"🎉 Presentation saved to: {output_path}")
    return True

def main():
    parser = argparse.ArgumentParser(description="Generate enhanced PowerPoint presentation")
    parser.add_argument("--schema", required=True, help="Path to the layout schema JSON file")
    parser.add_argument("--content", required=True, help="Path to the content JSON file")
    parser.add_argument("--output", default="output.pptx", help="Output PPTX file path")
    parser.add_argument("--images", default="assets", help="Directory containing images")
    
    args = parser.parse_args()
    
    build_enhanced_presentation(args.schema, args.content, args.output, args.images)

if __name__ == "__main__":
    main()